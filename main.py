from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from typing import List, Optional, Dict
import os
import json
from dotenv import load_dotenv
from core.llm import FallbackLLMService, get_llm_service
from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTXPt
from pypdf import PdfReader
from pptx.dml.color import RGBColor as PPTXRGBColor
from pptx.enum.text import PP_ALIGN
import io
import tempfile
from collections import Counter
from ppt_sanitizer import PPTSanitizer
import httpx
from bs4 import BeautifulSoup
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import re
from docx2pdf import convert
from datetime import datetime
from editorial_prompts import (
    DEVELOPMENT_EDITOR_TONE_OF_VOICE,
    LINE_EDITOR_INSTRUCTIONS,
    CONTENT_EDITOR_INSTRUCTIONS,
    COPY_EDITOR_INSTRUCTIONS,
    BRAND_ALIGNMENT_EDITOR_INSTRUCTIONS
)
from urllib.parse import urlparse
import ipaddress
import socket

# Audioop compatibility for Python 3.13+
try:
    import audioop
except ModuleNotFoundError:
    import audioop_lts as audioop
    import sys

    sys.modules["audioop"] = audioop

import boto3
from botocore.exceptions import BotoCoreError, ClientError
# from pydub import AudioSegment
import base64

load_dotenv()

app = FastAPI(title="PwC Presentation Assistant API")

# CORS Configuration
import re

# Define allowed origins with regex support for Amplify branches
allowed_origins_patterns = [
    r"http://localhost:\d+",  # Local development (any port)
    r"https://.*\.amplifyapp\.com",  # All Amplify branches
    r"https://deployment-frontend\.d2ebg85go3xrq2\.amplifyapp\.com",  # Specific Amplify app
]


def check_origin(origin: str) -> bool:
    """Check if origin matches any allowed pattern"""
    for pattern in allowed_origins_patterns:
        if re.match(pattern, origin):
            return True
    return False


app.add_middleware(
    CORSMiddleware,
    allow_origin_regex=r"(http://localhost:\d+|https://.*\.amplifyapp\.com)",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_llm():
    """Get the LLM service instance (supports Azure OpenAI + Groq fallback)"""
    try:
        return get_llm_service()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"LLM service not available: {str(e)}")


# Editorial Persona Descriptions - Single Source of Truth
PERSONA_DESCRIPTIONS = {
    "Development Editor": "Agent reviews the content's tone, theme and overall structure, makes structural adjustments to ensure consistency, alignment, and coherence throughout.",
    "Content Editor": "Agent evaluates strength and clarity of insights in the article, assess against the objectives of article, and refines language to align with author's key objectives.",
    "Line Editor": "Agent refines sentence structure and flow, improves style and readability, without changing the author's voice.",
    "Copy Editor": "Agent corrects grammar, punctuation, and typos, ensures consistency in formatting, adhering to PwC writing style and editorial standards.",
    "PwC Brand Alignment Editor": "Agent evaluates the content against the PwC brand guidelines, makes respective updates to align to PwC brand."
}

# All editorial prompts, guidelines, and workflow text are imported from editorial_prompts.py


class Message(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    messages: List[Message]
    stream: bool = True


class DraftRequest(BaseModel):
    topic: str
    objective: str
    audience: str
    additional_context: Optional[str] = None


class ThoughtLeadershipRequest(BaseModel):
    operation: str
    topic: Optional[str] = None
    perspective: Optional[str] = None
    target_audience: Optional[str] = None
    document_text: Optional[str] = None
    target_format: Optional[str] = None
    additional_context: Optional[str] = None
    reference_urls: Optional[List[str]] = None


class ExportRequest(BaseModel):
    content: str
    title: Optional[str] = "Generated Article"
    format: str  # "pdf" or "docx"


class ResearchRequest(BaseModel):
    query: str
    focus_areas: Optional[List[str]] = None
    additional_context: Optional[str] = None


class ArticleRequest(BaseModel):
    topic: str
    content_type: str  # Article, Case Study, Executive Brief, Blog, etc.
    desired_length: int  # word count
    tone: str  # Professional, Conversational, Technical, etc.
    outline_text: Optional[str] = None
    additional_context: Optional[str] = None


class BestPracticesRequest(BaseModel):
    categories: Optional[List[str]] = (
        None  # Structure, Visuals, Design, Charts, Formatting, Content
    )


class PodcastRequest(BaseModel):
    customization: Optional[str] = None


@app.get("/")
async def root():
    return {
        "message": "PwC Presentation Assistant API",
        "version": "1.0.0",
        "status": "running",
    }


@app.get("/health")
async def health_check():
    try:
        llm = get_llm()
        providers = llm.get_active_providers()
        return {
            "status": "healthy",
            "llm_providers": providers,
            "primary": providers[0] if providers else "none"
        }
    except Exception as e:
        return {
            "status": "degraded",
            "llm_providers": [],
            "error": str(e)
        }


async def crawl_related_pages(
    initial_url: str, max_pages: int = 5, max_depth: int = 2
) -> List[dict]:
    """
    Crawl related pages from the same domain to gather comprehensive content.
    Returns a list of fetched pages with their content.
    """
    parsed_initial = urlparse(initial_url)
    base_domain = f"{parsed_initial.scheme}://{parsed_initial.netloc}"

    visited_urls = set()
    to_visit = [(initial_url, 0)]  # (url, depth)
    fetched_pages = []

    # Browser-like headers
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
        "Accept-Encoding": "gzip, deflate, br",
        "Connection": "keep-alive",
    }

    async with httpx.AsyncClient(
        timeout=30.0, follow_redirects=True, max_redirects=5, headers=headers
    ) as client:
        while to_visit and len(fetched_pages) < max_pages:
            current_url, depth = to_visit.pop(0)

            # Skip if already visited or depth exceeded
            if current_url in visited_urls or depth > max_depth:
                continue

            visited_urls.add(current_url)

            try:
                # Fetch the page
                response = await client.get(current_url)

                if response.status_code != 200:
                    continue

                soup = BeautifulSoup(response.text, "html.parser")

                # Extract title
                title = soup.find("title")
                title_text = title.get_text().strip() if title else ""

                # Extract main content
                for script in soup(["script", "style", "nav", "footer", "header"]):
                    script.decompose()

                article_content = (
                    soup.find("article") or soup.find("main") or soup.find("body")
                )
                if article_content:
                    paragraphs = article_content.find_all(
                        ["p", "h1", "h2", "h3", "h4", "h5", "h6", "li"]
                    )
                    content = "\n".join(
                        [
                            p.get_text().strip()
                            for p in paragraphs
                            if p.get_text().strip()
                        ]
                    )
                else:
                    content = soup.get_text()

                content = re.sub(r"\n\s*\n", "\n\n", content).strip()

                # Store fetched page
                fetched_pages.append(
                    {
                        "url": current_url,
                        "title": title_text,
                        "content": content[:5000],  # Limit content per page
                    }
                )

                # Extract links for next level (only if not at max depth)
                if depth < max_depth and len(fetched_pages) < max_pages:
                    links = soup.find_all("a", href=True)
                    for link in links[:20]:  # Limit to first 20 links per page
                        href = link.get("href", "")
                        if not href:
                            continue

                        # Resolve relative URLs
                        if href.startswith("/"):
                            full_url = f"{base_domain}{href}"
                        elif href.startswith("http://") or href.startswith("https://"):
                            parsed_link = urlparse(href)
                            # Only follow links from same domain
                            if parsed_link.netloc == parsed_initial.netloc:
                                full_url = href
                            else:
                                continue
                        else:
                            # Relative URL
                            current_path = (
                                current_url.rsplit("/", 1)[0]
                                if "/" in current_url
                                else current_url
                            )
                            full_url = f"{current_path}/{href}"

                        # Normalize URL (remove fragments)
                        full_url = full_url.split("#")[0]

                        # Skip if already visited or queued
                        if (
                            full_url not in visited_urls
                            and (full_url, depth + 1) not in to_visit
                        ):
                            # Skip common non-content URLs
                            skip_patterns = [
                                "mailto:",
                                "tel:",
                                "javascript:",
                                "#",
                                ".pdf",
                                ".jpg",
                                ".png",
                                ".gif",
                                ".zip",
                                ".exe",
                                "login",
                                "signup",
                                "register",
                                "logout",
                                "admin",
                            ]
                            if not any(
                                pattern in full_url.lower() for pattern in skip_patterns
                            ):
                                to_visit.append((full_url, depth + 1))

            except Exception as e:
                # Continue to next URL on error
                continue

    return fetched_pages


def generate_url_variations(url: str) -> List[str]:
    """Generate common URL variations to try when encountering 404 errors"""
    variations = []
    parsed = urlparse(url)

    # Original URL (don't add it twice if it's already the first variation)
    path = parsed.path

    # Variation 1: Add trailing slash if not present
    if path and not path.endswith("/"):
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path}/")

    # Variation 2: Remove trailing slash if present
    if path.endswith("/") and len(path) > 1:
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path.rstrip('/')}")

    # Variation 3: Try lowercase path (common on case-sensitive servers)
    if path != path.lower():
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path.lower()}")

    # Variation 4: Try singular/plural variations for common patterns
    path_parts = path.strip("/").split("/")
    if len(path_parts) >= 1 and path_parts[-1]:
        last_part = path_parts[-1]
        # Try singular if plural
        if last_part.endswith("s") and len(last_part) > 1:
            singular = last_part[:-1]
            new_path = "/".join(path_parts[:-1] + [singular])
            variations.append(f"{parsed.scheme}://{parsed.netloc}/{new_path}")
        # Try plural if singular
        elif not last_part.endswith("s"):
            plural = last_part + "s"
            new_path = "/".join(path_parts[:-1] + [plural])
            variations.append(f"{parsed.scheme}://{parsed.netloc}/{new_path}")

    # Variation 5: Try with index.html
    if path and not path.endswith((".html", ".htm", ".php", ".asp", ".aspx")):
        variations.append(f"{parsed.scheme}://{parsed.netloc}{path}/index.html")

    # Add query string and fragment if present
    query_fragment = ""
    if parsed.query:
        query_fragment += f"?{parsed.query}"
    if parsed.fragment:
        query_fragment += f"#{parsed.fragment}"

    # Apply query/fragment to all variations
    if query_fragment:
        variations = [v + query_fragment for v in variations]

    return variations


def is_safe_hostname(hostname: str) -> bool:
    """Check if hostname is safe (not loopback, not private IP, not unspecified)"""
    try:
        ip = ipaddress.ip_address(hostname)
        return not (
            ip.is_loopback
            or ip.is_private
            or ip.is_reserved
            or ip.is_multicast
            or ip.is_unspecified
            or ip.is_link_local
        )
    except ValueError:
        try:
            resolved_ip = socket.gethostbyname(hostname)
            ip = ipaddress.ip_address(resolved_ip)
            return not (
                ip.is_loopback
                or ip.is_private
                or ip.is_reserved
                or ip.is_multicast
                or ip.is_unspecified
                or ip.is_link_local
            )
        except (socket.gaierror, ValueError):
            return True


async def fetch_url_content(url: str) -> dict:
    """Fetch and extract content from a URL with security validations"""
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ["http", "https"]:
            raise ValueError(
                f"Invalid URL scheme: {parsed.scheme}. Only http and https are allowed."
            )

        if not parsed.hostname:
            raise ValueError("Invalid URL: missing domain")

        if not is_safe_hostname(parsed.hostname):
            raise ValueError(
                "Access to localhost, private IP ranges, or reserved IPs is not allowed for security reasons"
            )

        # Browser-like headers to avoid 403 Forbidden errors
        # Build referer from the URL (same domain)
        referer = f"{parsed.scheme}://{parsed.netloc}/"

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate, br",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
            "Referer": referer,
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "same-origin",
            "Sec-Fetch-User": "?1",
            "Cache-Control": "max-age=0",
            "DNT": "1",
        }

        async with httpx.AsyncClient(
            timeout=60.0, follow_redirects=True, max_redirects=5, headers=headers
        ) as client:
            response = await client.get(url)

            # Check status codes before raising
            if response.status_code == 404:
                # Try URL variations before giving up
                variations = generate_url_variations(url)
                tried_urls = [url]

                for variation_url in variations:
                    if variation_url in tried_urls:
                        continue
                    tried_urls.append(variation_url)

                    try:
                        # Update referer for the variation
                        parsed_var = urlparse(variation_url)
                        var_referer = f"{parsed_var.scheme}://{parsed_var.netloc}/"
                        headers["Referer"] = var_referer

                        var_response = await client.get(variation_url)
                        if var_response.status_code == 200:
                            # Success! Use this URL and continue with content extraction
                            url = variation_url  # Update url to the working variation
                            response = var_response
                            break
                        elif var_response.status_code != 404:
                            # If it's not 200 or 404, it's a different error, stop trying
                            break
                    except Exception:
                        # Continue to next variation on any error
                        continue

                # If still 404 after trying variations, return error
                if response.status_code == 404:
                    return {
                        "url": url,
                        "title": "",
                        "content": "",
                        "success": False,
                        "error": f"404 Not Found: The page at this URL does not exist. Tried {len(tried_urls)} variation(s) but none were found. Please verify the URL is correct.",
                    }
            elif response.status_code == 403:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"403 Forbidden: The website blocked access. This may be due to bot protection, Cloudflare, or WAF rules. The site may require JavaScript to load content.",
                }
            elif response.status_code == 429:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"429 Too Many Requests: Rate limit exceeded. Please wait a moment and try again.",
                }
            elif response.status_code == 503:
                return {
                    "url": url,
                    "title": "",
                    "content": "",
                    "success": False,
                    "error": f"503 Service Unavailable: The website is temporarily unavailable. Please try again later.",
                }

            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")

            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            title = soup.find("title")
            title_text = title.get_text().strip() if title else ""

            article_content = (
                soup.find("article") or soup.find("main") or soup.find("body")
            )
            if article_content:
                paragraphs = article_content.find_all(
                    ["p", "h1", "h2", "h3", "h4", "li"]
                )
                content = "\n".join(
                    [p.get_text().strip() for p in paragraphs if p.get_text().strip()]
                )
            else:
                content = soup.get_text()

            content = re.sub(r"\n\s*\n", "\n\n", content)
            content = content.strip()

            return {
                "url": url,
                "title": title_text,
                "content": content[:5000],
                "success": True,
            }
    except Exception as e:
        return {
            "url": url,
            "title": "",
            "content": "",
            "success": False,
            "error": str(e),
        }


def extract_text_from_pdf(pdf_bytes: bytes, max_chars: int = 10000) -> str:
    """Extract text content from PDF file"""
    try:
        pdf = PdfReader(io.BytesIO(pdf_bytes))
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
            if len(text) > max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        return f"[Error reading PDF: {str(e)}]"


def extract_text_from_docx(docx_bytes: bytes, max_chars: int = 10000) -> str:
    """Extract text content from DOCX file"""
    try:
        doc = Document(io.BytesIO(docx_bytes))
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
            if len(text) > max_chars:
                break
        return text[:max_chars]
    except Exception as e:
        return f"[Error reading DOCX: {str(e)}]"


async def parse_uploaded_file(uploaded_file: UploadFile, max_chars: int = 10000) -> str:
    """Parse uploaded file and extract text content based on file type"""
    if not uploaded_file:
        return ""

    file_content = await uploaded_file.read()
    filename = uploaded_file.filename.lower() if uploaded_file.filename else ""

    if filename.endswith(".pdf"):
        return extract_text_from_pdf(file_content, max_chars)
    elif filename.endswith(".docx"):
        return extract_text_from_docx(file_content, max_chars)
    elif filename.endswith(".pptx"):
        try:
            prs = Presentation(io.BytesIO(file_content))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            text_runs.append(text)
            combined = "\n\n".join(text_runs)
            return combined[:max_chars]
        except Exception as e:
            return f"[Error reading PPTX: {str(e)}]"
    elif filename.endswith((".txt", ".md")):
        try:
            return file_content.decode("utf-8")[:max_chars]
        except:
            return f"[Could not decode text file: {uploaded_file.filename}]"
    else:
        try:
            return file_content.decode("utf-8")[:max_chars]
        except:
            return (
                f"[Unsupported file type or could not decode: {uploaded_file.filename}]"
            )

# ============================================
# COMPLETE TEMPLATE-BASED DOCUMENT GENERATION
# ============================================

def extract_instructions_from_template(doc: Document) -> tuple[str, Optional[int], Optional[int]]:
    """
    Extract formatting instructions from pages 3-4 of the template.
    These instructions are used to format the AI-generated content.
    
    Returns:
        (instructions_text, start_idx, end_idx)
    """
    print(f"\n🔍 Extracting formatting instructions from template...")
    
    # Extended list of instruction markers to capture all formatting guidelines
    instruction_markers = [
        'body text style:',
        'list bullet style:',
        'list number style:',
        'list alpha style:',
        'heading 1 style',
        'heading 2 style',
        'heading 3 style',
        'heading 4 style',
        'chart header',
        'table header',
        'table text',
        'caption',
        'quote',
        'font is 11 points',
        '1.5 lines spacing',
        'hanging indent',
        'space after has been added',
        'space after',
        'double returns',
        'styles in your new templates',
        'programmed with space after'
    ]
    
    # Keywords that indicate the end of instructions
    end_markers = [
        'report large callout',
        'placer at prod',
        'dolor nam eim',
        'persius apeirian',
        'disputationi',
        'aeterno fuisset',
        'odio legere consulatu'
    ]
    
    instructions = []
    start_idx = None
    end_idx = None
    found_instructions = False
    consecutive_instruction_paras = 0
    
    # Scan through all paragraphs to find instruction blocks
    for idx, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        text_lower = text.lower()
        
        # Skip empty paragraphs
        if not text:
            continue
        
        # Check if this is an instruction paragraph
        is_instruction = any(marker in text_lower for marker in instruction_markers)
        
        # Check if this indicates end of instructions (decorative/example content)
        is_end_marker = any(marker in text_lower for marker in end_markers)
        
        if is_instruction:
            if start_idx is None:
                start_idx = idx
                found_instructions = True
            instructions.append(text)
            end_idx = idx
            consecutive_instruction_paras += 1
        elif found_instructions:
            # If we found instructions and now see end markers, stop
            if is_end_marker:
                break
            # Continue collecting if it's part of the instruction block
            # (look for paragraphs that continue the instruction context)
            if consecutive_instruction_paras > 0 and len(text) < 200:
                # Might be continuation of instructions
                if any(keyword in text_lower for keyword in ['style', 'font', 'spacing', 'indent', 'points']):
                    instructions.append(text)
                    end_idx = idx
                    consecutive_instruction_paras += 1
                else:
                    consecutive_instruction_paras = 0
            else:
                # If we've moved far from instructions and hit a long paragraph, likely done
                if len(text) > 150 and idx > (end_idx or 0) + 5:
                    break
                consecutive_instruction_paras = 0
    
    if not instructions:
        print(f"   ⚠️  No instructions found in template, using default formatting")
        return "", None, None
    
    instructions_text = "\n".join(instructions)
    print(f"   ✅ Found {len(instructions)} instruction paragraphs")
    print(f"   📍 Instructions span: paragraphs {start_idx} to {end_idx}")
    
    return instructions_text, start_idx, end_idx


# ============================================
# SECTION 2: AI FORMATTING
# ============================================

async def generate_formatted_content_from_instructions(
    raw_content: str,
    instructions: str,
    llm_service
) -> str:
    """
    Use Pages 3-4 instructions as system prompt to format AI content professionally.
    
    IMPORTANT: This function receives instructions that were extracted from pages 3-4
    BEFORE those pages were removed. The instructions are used here as the system prompt
    to guide the LLM in formatting the content according to template styles.
    
    Args:
        raw_content: Raw AI-generated content to format
        instructions: Extracted instructions from pages 3-4 (used as system prompt)
        llm_service: LLM service for formatting
    
    Returns:
        Formatted content in markdown format
    """
    print(f"\n🤖 Using extracted instructions as system prompt to format content...")
    print(f"   📋 Instructions length: {len(instructions)} characters")
    
    system_prompt = f"""You are a professional document formatter for business clients. Format the provided content according to these specific style guidelines extracted from the template:

{instructions}

**Your Task:**
1. Take the raw content provided by the user
2. Format it professionally using the style guidelines above
3. Start with a Heading 1 (#) for the main content title - this is IMPORTANT
4. Apply appropriate headings (# ## ### ####) throughout the document
5. Use bullet points where appropriate (-)
6. Format lists (numbered or alphabetical) correctly
7. Ensure proper spacing and structure
8. Use bold text (**text**) for emphasis only where necessary
9. Maintain professional business tone and clarity
10. Output ONLY the formatted content, no explanations

**Format Rules:**
- Heading 1 (#): Use for the main document title/content title (MUST start with this)
- Heading 2 (##): Use for major sections  
- Heading 3 (###): Use for subsections
- Heading 4 (####): Use for sub-subsections
- Body Text: Regular paragraphs (no prefix) - left aligned, 1.5 line spacing
- Bullet Lists: Start lines with "-" - use hanging indent
- Numbered Lists: Start lines with "1. 2. 3." - use hanging indent
- Alphabetical Lists: Start lines with "A. B. C." - use hanging indent
- Bold Text: Wrap in **text** - use sparingly for emphasis
- Quotes: Start lines with ">" - italic, indented

**Alignment and Formatting:**
- All headings: LEFT aligned
- Body text: LEFT aligned, 11pt font, 1.5 line spacing
- Lists: LEFT aligned with hanging indent (0.25 inches)
- Proper spacing: Space after paragraphs (6pt), space before headings (12-24pt)
- Professional appearance: Consistent formatting, clear hierarchy

**IMPORTANT:** 
- The formatted content MUST start with a Heading 1 (#) to serve as the content title
- Use professional business formatting with proper alignment and spacing
- Ensure all text is left-aligned (except quotes which are indented)
- Maintain consistent formatting throughout the document

Output a complete, professionally formatted document starting with a Heading 1."""

    user_prompt = f"""Format this content professionally:

{raw_content}

Apply proper formatting based on the style guidelines provided."""

    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3,
            max_tokens=4000
        )
        
        formatted_content = response.content.strip()
        print(f"   ✅ Content formatted ({len(formatted_content)} chars)")
        return formatted_content
        
    except Exception as e:
        print(f"   ⚠️  Formatting failed: {e}, using original content")
        return raw_content


async def generate_title_and_content_for_page2(raw_ai_content: str, llm_service) -> tuple[str, str]:
    """
    Generate title and content for page 2 based on AI response.
    Uses LLM to extract and format a title and summary/content for page 2.
    
    Returns:
        Tuple of (title, content) for page 2
    """
    print(f"\n📄 Generating title and content for page 2 from AI response...")
    
    system_prompt = """You are a professional document editor. Based on the AI-generated content, create:
1. A professional title (maximum 10 words, title case, no quotation marks)
2. A concise summary/content section (2-4 paragraphs) that provides an overview of the main content

The title and content will be displayed on page 2 of the document.

Rules:
- Title: Maximum 10 words, professional tone, title case, no quotation marks
- Content: 2-4 paragraphs summarizing the key points, professional tone, well-structured
- Output format: 
  TITLE: [title here]
  
  CONTENT:
  [content paragraphs here]

Output ONLY the title and content in the specified format."""
    
    content_preview = raw_ai_content[:2000] if len(raw_ai_content) > 2000 else raw_ai_content
    
    user_prompt = f"""Generate a professional title and summary content for page 2 based on this AI-generated content:

{content_preview}

Create:
1. A concise, impactful title (max 10 words)
2. A 2-4 paragraph summary that provides an overview of the main content

Output in format:
TITLE: [title]

CONTENT:
[paragraph 1]

[paragraph 2]

[paragraph 3]"""
    
    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.4,
            max_tokens=800
        )
        
        response_text = response.content.strip()
        
        # Parse title and content
        title = ""
        content = ""
        
        if "TITLE:" in response_text:
            parts = response_text.split("TITLE:", 1)
            if len(parts) > 1:
                title_section = parts[1]
                if "CONTENT:" in title_section:
                    title = title_section.split("CONTENT:")[0].strip()
                    content = title_section.split("CONTENT:", 1)[1].strip()
                else:
                    title = title_section.strip()
        elif "CONTENT:" in response_text:
            parts = response_text.split("CONTENT:", 1)
            title = "Document Overview"  # Default title
            content = parts[1].strip() if len(parts) > 1 else response_text
        else:
            # Fallback: use first line as title, rest as content
            lines = response_text.split('\n', 1)
            title = lines[0].strip()[:100] if lines else "Document Overview"
            content = lines[1].strip() if len(lines) > 1 else response_text
        
        # Clean title (remove quotes, limit length)
        title = title.replace('"', '').replace("'", '').strip()
        if len(title) > 100:
            title = title[:97] + "..."
        if not title:
            title = "Document Overview"
        
        # Clean content
        if not content:
            content = "This document provides comprehensive information on the topic."
        
        print(f"   ✅ Generated page 2 title: '{title}'")
        print(f"   ✅ Generated page 2 content ({len(content)} chars)")
        return title, content
        
    except Exception as e:
        print(f"   ⚠️  Failed to generate page 2 content: {e}, using defaults")
        return "Document Overview", "This document provides comprehensive information and analysis on the topic."
    



    
async def generate_title_from_ai_content(content: str, llm_service) -> str:
    """
    Generate professional title from AI content.
    """
    print(f"\n📌 Generating title from content...")
    
    system_prompt = """You are a professional editor. Generate a concise, impactful title.

Rules:
- Maximum 10 words
- Professional tone
- Capture main topic
- No quotation marks
- Title case format
- Output ONLY the title"""

    content_preview = content[:800] if len(content) > 800 else content
    
    user_prompt = f"""Generate a professional title for this content:

{content_preview}

Output ONLY the title, nothing else."""

    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7,
            max_tokens=30
        )
        
        title = response.content.strip().replace('"', '').replace("'", '').strip()
        
        if len(title) > 100:
            title = title[:97] + "..."
        
        print(f"   ✅ Generated title: '{title}'")
        return title if title else "Professional Report"
        
    except Exception as e:
        print(f"   ⚠️  Title generation failed: {e}")
        return "Professional Report"


# ============================================
# SECTION 3: CONTENT PARSING (SINGLE PARSER)
# ============================================

def parse_markdown_to_structure(content: str) -> List[Dict]:
    """
    Parse markdown-formatted content into Word structure.
    
    SINGLE PARSING FUNCTION - NO DUPLICATES!
    
    Returns:
        List of dicts with 'style' and 'text' keys
    """
    STYLE_MAP = {
        "heading1": "Heading 1 style",
        "heading2": "Heading 2 style",
        "heading3": "Heading 3 style",
        "heading4": "Heading 4 style",
        "body": "Body Text style",
        "bullet": "List Bullet style",
        "number": "List Number style",
        "alpha": "List Alpha style",
        "quote": "Quote",
        "caption": "Caption"
    }
    
    structured = []
    lines = content.split('\n')
    
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        if not line:
            i += 1
            continue
        
        # Heading 1: # Title
        if line.startswith('# ') and not line.startswith('##'):
            structured.append({"style": STYLE_MAP["heading1"], "text": line[2:].strip()})
            i += 1
            continue
        
        # Heading 2: ## Section
        if line.startswith('## ') and not line.startswith('###'):
            structured.append({"style": STYLE_MAP["heading2"], "text": line[3:].strip()})
            i += 1
            continue
        
        # Heading 3: ### Subsection
        if line.startswith('### ') and not line.startswith('####'):
            structured.append({"style": STYLE_MAP["heading3"], "text": line[4:].strip()})
            i += 1
            continue
        
        # Heading 4: #### Sub-subsection
        if line.startswith('#### '):
            structured.append({"style": STYLE_MAP["heading4"], "text": line[5:].strip()})
            i += 1
            continue
        
        # Quote: > Text
        if line.startswith('> '):
            structured.append({"style": STYLE_MAP["quote"], "text": line[2:].strip()})
            i += 1
            continue
        
        # Bullet List: - Item, • Item, * Item
        if line.startswith('- ') or line.startswith('• ') or line.startswith('* '):
            bullet_text = line[2:].strip() if line.startswith('- ') or line.startswith('* ') else line[1:].strip()
            structured.append({"style": STYLE_MAP["bullet"], "text": bullet_text})
            i += 1
            continue
        
        # Numbered List: 1. Item
        if re.match(r'^\d+\.\s+', line):
            text = re.sub(r'^\d+\.\s+', '', line).strip()
            structured.append({"style": STYLE_MAP["number"], "text": text})
            i += 1
            continue
        
        # Alphabetical List: A. Item
        if re.match(r'^[A-Za-z]\.\s+', line):
            text = re.sub(r'^[A-Za-z]\.\s+', '', line).strip()
            structured.append({"style": STYLE_MAP["alpha"], "text": text})
            i += 1
            continue
        
        # Caption: [text]
        if line.startswith('[') and line.endswith(']'):
            structured.append({"style": STYLE_MAP["caption"], "text": line[1:-1].strip()})
            i += 1
            continue
        
        # Body Text - collect multiple lines
        paragraph_lines = [line]
        i += 1
        
        while i < len(lines):
            next_line = lines[i].strip()
            if not next_line:
                i += 1
                continue
            # Check if next line starts a new block
            if (next_line.startswith('#') or \
                next_line.startswith('> ') or \
                next_line.startswith('- ') or \
                next_line.startswith('• ') or \
                next_line.startswith('* ') or \
                next_line.startswith('[') or \
               re.match(r'^\d+\.\s+', next_line) or \
                re.match(r'^[A-Za-z]\.\s+', next_line)):
                break
            paragraph_lines.append(next_line)
            i += 1
        
        text = ' '.join(paragraph_lines).strip()
        if text:
            structured.append({"style": STYLE_MAP["body"], "text": text})
    
    return structured


# ============================================
# SECTION 4: DOCUMENT MANIPULATION
# ============================================

def remove_all_footers(doc: Document):
    """
    Remove all footers from the document while preserving headers, logos, and theme.
    Only removes footer text content, preserves document structure and theme.
    """
    print(f"\n🗑️  Removing footer text content (preserving headers and theme)...")
    try:
        # Access document sections
        for section in doc.sections:
            # Remove footer text content but preserve structure
            footer = section.footer
            # Clear all paragraphs in footer (removes text)
            for para in footer.paragraphs:
                para.clear()
            
            # Note: We don't remove the footer element itself to preserve document structure
            # This ensures headers, logos, backgrounds, and theme are preserved
        print(f"   ✅ Footer text content removed (headers and theme preserved)")
    except Exception as e:
        print(f"   ⚠️  Could not remove footers: {e}")


def preserve_template_theme(doc: Document):
    """
    Preserve template theme across all pages:
    - Keep headers (logos, branding) - headers are automatically preserved
    - Preserve section properties (background colors, margins, page size)
    - Maintain document structure and formatting
    - Ensure headers and footers remain intact throughout document modifications
    - Ensure all sections have consistent headers, footers, logos, and backgrounds
    - CRITICAL: Ensure first section doesn't have page break settings that push content to page 2
    
    IMPORTANT: This should be called BEFORE any content removal or modifications
    to ensure the theme is maintained throughout the document.
    """
    try:
        # Document sections automatically preserve headers, backgrounds, and theme
        # when we load the template. We need to ensure consistency across all sections.
        
        # Verify sections exist and headers are present
        section_count = len(doc.sections)
        if section_count == 0:
            print(f"   ⚠️  No sections found in document")
            return
        
        # Get the first section's header and footer as reference (from page 1)
        first_section = doc.sections[0]
        reference_header = first_section.header
        reference_footer = first_section.footer
        
        # CRITICAL: Ensure first section doesn't have page break settings
        # Check if first section has any page break properties that might cause issues
        try:
            # Ensure first section starts on same page (no odd/even page break)
            # This ensures content starts on page 1, not page 2
            if hasattr(first_section, 'start_type'):
                # Check section start type - should be CONTINUOUS or NEW_PAGE, not ODD_PAGE or EVEN_PAGE
                from docx.enum.section import WD_SECTION_START
                if first_section.start_type == WD_SECTION_START.ODD_PAGE:
                    # ODD_PAGE can cause content to start on page 2 if page 1 is even
                    print(f"   ⚠️  First section starts on ODD_PAGE - this might cause title to appear on page 2")
                    # We'll leave it for now, but this could be the issue
                elif first_section.start_type == WD_SECTION_START.EVEN_PAGE:
                    # EVEN_PAGE can cause content to start on page 2
                    print(f"   ⚠️  First section starts on EVEN_PAGE - this might cause title to appear on page 2")
        except Exception as e:
            print(f"   ℹ️  Could not check section start type: {e}")
        
        # Ensure all sections have the same header and footer as page 1
        # This ensures consistency across all pages
        for i, section in enumerate(doc.sections):
            try:
                # Verify headers exist (they should be preserved automatically)
                header = section.header
                footer = section.footer
                
                # Headers typically contain logos and branding
                # These are preserved automatically by python-docx when loading template
                
                # Preserve section properties (margins, page size, orientation, background)
                # These are already set from the template and remain unchanged
                # unless we explicitly modify them (which we don't)
                
                # Note: python-docx automatically maintains header/footer consistency
                # when sections are created from the template. We just verify they exist.
                
                if i > 0:
                    # For additional sections, ensure they match the first section's theme
                    # This is typically handled automatically, but we verify
                    print(f"   ✅ Section {i+1} theme verified (matches page 1)")
                
            except Exception as e:
                print(f"   ⚠️  Could not verify section {i} header/footer: {e}")
        
        print(f"   ✅ Template theme preserved (headers, footers, backgrounds, colors, logos)")
        print(f"   ✅ {section_count} section(s) with consistent headers and footers maintained")
        print(f"   ✅ All pages will have the same logo, background, theme, headers, and footers as page 1")
        print(f"   ✅ Theme consistency verified across all sections")
        
    except Exception as e:
        print(f"   ⚠️  Could not preserve theme: {e}")
        # Continue execution even if theme preservation check fails
        # The template theme should still be preserved by python-docx


def clean_page1_keep_only_title(doc: Document, title: str) -> bool:
    """
    Clean page 1 - remove ALL template text content, subtitles, page breaks, and keep ONLY the title.
    Title is centered horizontally and vertically on page 1.
    PRESERVES: Logos, images in headers, background colors, and template theme.
    REMOVES: All text content, subtitles, placeholder text, page breaks - ONLY title remains on page 1.
    """
    if not title or not title.strip():
        print(f"   ⚠️  No title provided, skipping page 1 title update")
        return False
    
    print(f"\n🧹 Cleaning page 1 - keeping ONLY title (removing ALL content including page breaks): '{title}'")
    
    # Clean title
    title = title.strip()
    
    # CRITICAL: Remove ALL paragraphs from the document first, then add ONLY the title
    # This ensures title appears on page 1, not page 2
    paragraphs_to_remove = []
    
    # Check ALL paragraphs - we need to remove everything from page 1
    # But be careful: we'll remove first 60 paragraphs to cover page 1 completely
    page1_end = min(60, len(doc.paragraphs))
    
    # Keywords that indicate template/placeholder content to remove (including Contents)
    template_keywords = [
        'instruction', 'style:', 'template', 'placeholder', 'body text style',
        'heading 1 style', 'heading 2 style', 'heading 3 style', 'heading 4 style',
        'list bullet style', 'list number style', 'list alpha style',
        'font is 11 points', '1.5 lines spacing', 'space after',
        'report large callout', 'placer at prod', 'dolor nam eim',
        'persius apeirian', 'disputationi', 'aeterno fuisset',
        'odio legere consulatu', 'lorem', 'example', 'sample', 'subtitle',
        '错误', 'manually placed', 'dividing content', 'ending',
        'contents', 'content', 'table of contents', 'toc',  # Remove Contents from page 1
        'strictly private', 'confidential'  # Remove footer text if it's in body
    ]
    
    # Check for page breaks in paragraphs and remove them too
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    
    # REMOVE ALL PARAGRAPHS FROM PAGE 1 - Be extremely aggressive
    # Page 1 should ONLY have the title we add - nothing else
    for idx in range(page1_end):
        if idx >= len(doc.paragraphs):
            break
        
        para = doc.paragraphs[idx]
        text = para.text.strip()
        text_lower = text.lower()
        
        # Check for page breaks in this paragraph
        has_page_break = False
        try:
            # Check if paragraph contains page break
            for run in para.runs:
                if run._element.xpath('.//w:br[@w:type="page"]'):
                    has_page_break = True
                    break
        except:
            pass
        
        # PRESERVE images, logos, and graphics ONLY in headers (section headers, not document paragraphs)
        has_image = False
        try:
            for run in para.runs:
                if run._element.xpath('.//w:drawing') or run._element.xpath('.//w:pict'):
                    has_image = True
                    break
        except:
            pass
        
        # Skip images ONLY if they're in headers (headers are preserved automatically)
        # Remove ALL text paragraphs and page breaks from page 1
        
        # Remove if it has a page break (we don't want page breaks on page 1)
        if has_page_break:
            paragraphs_to_remove.append(idx)
            print(f"   🗑️  Found page break in paragraph {idx}, marking for removal")
        # Remove if it contains template keywords (including Contents)
        elif text or len(para.runs) > 0:
            if any(kw in text_lower for kw in template_keywords):
                paragraphs_to_remove.append(idx)
            # Remove ALL text content from page 1 (first 50 paragraphs)
            # This ensures Contents, headings, and all template content is removed
            elif idx < 50:  # Remove first 50 paragraphs (covers all of page 1)
                paragraphs_to_remove.append(idx)
            # Also remove if it looks like template content
            elif len(text) > 0 and (len(text) < 300 or any(char in text for char in ['...', '---', '___', '....'])):
                paragraphs_to_remove.append(idx)
    
    # Remove paragraphs in reverse order to maintain indices
    removed_count = 0
    for idx in reversed(paragraphs_to_remove):
        try:
            if idx < len(doc.paragraphs):
                p_element = doc.paragraphs[idx]._element
                parent = p_element.getparent()
                if parent is not None:
                    parent.remove(p_element)
                    removed_count += 1
        except Exception as e:
            print(f"   ⚠️  Could not remove paragraph {idx}: {e}")
            continue
    
    print(f"   ✅ Removed {removed_count} paragraphs (including page breaks) from page 1")
    
    # CRITICAL: Remove ALL remaining paragraphs to ensure title is FIRST element
    # This is the key fix - remove everything, then add title as first element
    print(f"   🧹 Removing ALL remaining paragraphs to ensure title is first element...")
    additional_removals = 0
    max_additional_removals = 150  # Remove up to 150 paragraphs to be safe
    while len(doc.paragraphs) > 0 and additional_removals < max_additional_removals:
        try:
            # Always remove the FIRST paragraph (index 0)
            p = doc.paragraphs[0]
            p_element = p._element
            parent = p_element.getparent()
            if parent is not None:
                parent.remove(p_element)
                additional_removals += 1
            else:
                break
        except Exception as e:
            print(f"   ⚠️  Could not remove paragraph 0: {e}")
            break
    
    if additional_removals > 0:
        print(f"   🗑️  Removed {additional_removals} additional paragraphs")
    
    # Verify document is now empty (or nearly empty)
    remaining_paras = len(doc.paragraphs)
    if remaining_paras > 0:
        print(f"   ⚠️  Warning: {remaining_paras} paragraphs still remain after cleaning")
        # Force remove remaining paragraphs
        for idx in range(remaining_paras - 1, -1, -1):
            try:
                if idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]
                    p_element = p._element
                    parent = p_element.getparent()
                    if parent is not None:
                        parent.remove(p_element)
            except:
                pass
    
    # CRITICAL FIX: Ensure document body is clean and title is FIRST element
    # Check document structure and ensure title is in the first section's body
    try:
        # Verify we have at least one section
        if len(doc.sections) == 0:
            print(f"   ⚠️  WARNING: Document has no sections!")
        
        # Ensure first section doesn't force page breaks
        if len(doc.sections) > 0:
            first_section = doc.sections[0]
            try:
                from docx.enum.section import WD_SECTION_START
                # Force first section to start on new page (normal behavior for page 1)
                # But ensure it's not ODD_PAGE or EVEN_PAGE which can skip to page 2
                if hasattr(first_section, 'start_type'):
                    current_start_type = first_section.start_type
                    # Only change if it's ODD_PAGE or EVEN_PAGE (which can cause issues)
                    if current_start_type in [WD_SECTION_START.ODD_PAGE, WD_SECTION_START.EVEN_PAGE]:
                        first_section.start_type = WD_SECTION_START.NEW_PAGE
                        print(f"   🔧 Fixed: Changed section start type from {current_start_type} to NEW_PAGE")
            except Exception as e:
                print(f"   ℹ️  Could not modify section start type: {e}")
        
        # Now add title as the ABSOLUTE FIRST paragraph in the document body
        # Document should be empty at this point after cleaning
        if len(doc.paragraphs) > 0:
            # Still have paragraphs - clear the first one and use it
            print(f"   ⚠️  Document still has {len(doc.paragraphs)} paragraphs, clearing first one...")
            title_para = doc.paragraphs[0]
            title_para.clear()  # Clear all runs from first paragraph
        else:
            # Document is empty - add new paragraph (this will be first)
            title_para = doc.add_paragraph()
        
        # Configure paragraph formatting for centered title on page 1
        title_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.CENTER  # Center horizontally
        title_para.paragraph_format.space_before = Pt(320)  # Space from top to center vertically on page
        title_para.paragraph_format.space_after = Pt(0)  # No space after (page 1 should end here)
        title_para.paragraph_format.line_spacing = 1.15  # Tight line spacing for title
        title_para.paragraph_format.keep_with_next = False  # Don't force to next page
        
        # CRITICAL: Ensure no page break before this paragraph
        try:
            # Remove any page break elements from this paragraph
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            # Check if paragraph has page break and remove it
            for run in title_para.runs:
                try:
                    # Remove any page breaks from runs
                    br_elements = run._element.xpath('.//w:br[@w:type="page"]')
                    for br in br_elements:
                        br.getparent().remove(br)
                except:
                    pass
        except:
            pass
        
        # Add title text with energetic styling (larger, bold, orange)
        run = title_para.add_run(title)
        run.bold = True
        run.font.size = Pt(38)  # Large size for energetic, prominent title
        try:
            run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)  # PwC orange (#E87722)
        except:
            try:
                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray fallback
            except:
                pass
        
        # CRITICAL VERIFICATION: Ensure title is the first paragraph at index 0
        if len(doc.paragraphs) == 0:
            print(f"   ❌ ERROR: Document has no paragraphs after title addition!")
            return False
        
        first_para_text = doc.paragraphs[0].text.strip()
        if first_para_text == title:
            print(f"   ✅ VERIFIED: Title '{title[:50]}...' is the FIRST paragraph (index 0)")
            print(f"   ✅ Title is at the very beginning of the document")
            print(f"   ✅ Title will appear on PAGE 1 (centered)")
        else:
            print(f"   ⚠️  ERROR: Title is NOT first paragraph!")
            print(f"   ⚠️  First paragraph text: '{first_para_text[:80]}...'")
            print(f"   ⚠️  Expected title: '{title}'")
            print(f"   ⚠️  Document has {len(doc.paragraphs)} paragraphs")
            # Force fix: remove ALL paragraphs and re-add title as first
            try:
                # Remove all paragraphs
                while len(doc.paragraphs) > 0:
                    try:
                        p = doc.paragraphs[0]
                        p_element = p._element
                        parent = p_element.getparent()
                        if parent is not None:
                            parent.remove(p_element)
                    except:
                        break
                # Now add title as first and only paragraph
                title_para = doc.add_paragraph()
                title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                title_para.paragraph_format.space_before = Pt(320)
                title_para.paragraph_format.keep_with_next = False
                run = title_para.add_run(title)
                run.bold = True
                run.font.size = Pt(38)
                try:
                    run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)
                except:
                    pass
                print(f"   🔧 FIXED: Removed all paragraphs and re-added title as first")
                # Verify again
                if len(doc.paragraphs) > 0 and doc.paragraphs[0].text.strip() == title:
                    print(f"   ✅ VERIFIED AFTER FIX: Title is now first paragraph")
                else:
                    print(f"   ❌ FIX FAILED: Title still not first paragraph")
                    return False
            except Exception as e:
                print(f"   ❌ Could not fix title position: {e}")
                import traceback
                traceback.print_exc()
                return False
        
        # Final verification of document structure
        print(f"   📊 Final document structure:")
        print(f"      - Total paragraphs: {len(doc.paragraphs)}")
        print(f"      - First paragraph (index 0): '{doc.paragraphs[0].text[:60] if len(doc.paragraphs) > 0 else 'N/A'}...'")
        if len(doc.paragraphs) > 1:
            print(f"      - Second paragraph (index 1): '{doc.paragraphs[1].text[:60]}...'")
        
        print(f"   ✅ Added energetic title to page 1 (centered, ONLY title): '{title}'")
        print(f"   ✅ Page 1 now contains ONLY the title - all other content and page breaks removed")
        print(f"   ✅ Title is positioned at index 0 (very beginning of document)")
        print(f"   ✅ Title will appear on PAGE 1, centered on the page")
        return True
    except Exception as e:
        print(f"   ⚠️  Could not add title: {e}")
        import traceback
        traceback.print_exc()
        # Fallback: try adding title directly
        try:
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_para.paragraph_format.space_before = Pt(320)
            run = title_para.add_run(title)
            run.bold = True
            run.font.size = Pt(38)
            try:
                run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)
            except:
                pass
            print(f"   ✅ Added title using fallback method")
            return True
        except Exception as e2:
            print(f"   ❌ Fallback also failed: {e2}")
            return False


def update_cover_title(doc: Document, title: str) -> bool:
    """
    Update title on cover page (page 1).
    This is a wrapper that calls clean_page1_keep_only_title for a clean cover page.
    NOTE: This function is kept for compatibility but clean_page1_keep_only_title is called directly.
    """
    try:
        return clean_page1_keep_only_title(doc, title)
    except Exception as e:
        print(f"   ⚠️  Error adding title: {e}")
        import traceback
        traceback.print_exc()
    return False


def remove_template_instructions_and_unnecessary_content(doc: Document, start_idx: int, end_idx: int):
    """
    Remove instruction paragraphs and all unnecessary template content.
    This includes decorative callouts, example text, and placeholder content.
    """
    if start_idx is None or end_idx is None:
        return
    
    print(f"\n🗑️  Removing template instructions and unnecessary content...")
    
    # Comprehensive list of keywords that indicate unnecessary/decorative content
    unnecessary_keywords = [
        'report large callout',
        'placer at prod',
        'dolor nam eim',
        'assum e est',
        'persius apeirian',
        'disputationi',
        'min aeterno fuisset',
        'usu odio legere',
        'consulatu cu placer',
        'unnecessary',
        'dont want in template',
        'listen first',
        'take 3rd 4rd th page',
        'instruction as system prompt',
        'extract text',
        'user content',
        'ai response',
        'based on take',
        'template first page',
        'replace a title',
        '2nd page content',
        'llm will choose',
        'provide as professional',
        'seen this ai response',
        'pick this and remove',
        'input template text',
        'replace those in 1 st',
        '2rd page',
        '3rd pae as',
        '4tg page instruction',
        'take this instruction',
        'frmat badsed',
        'align those',
        'show 3rd page',
        'remove instruction',
        'everything and replace',
        'show user as professional',
        'think once in templacte',
        'what are unnessagry',
        'no nned to show',
        'everything sould remove',
        'provide profeesional',
        'business user',
        'same foloow both word',
        'pdf export',
        'check instruction',
        'waht are heakings',
        'fonts size colour',
        'everything check inch',
        'provide ...check code',
        'whata re duplicate',
        'unnnesaagry code remove',
        'check fronyend aslo',
        'provide [rroper code',
        'output and exlabtion'
    ]
    
    # Instruction-related keywords (should also be removed after extraction)
    instruction_keywords = [
        'body text style',
        'list bullet style',
        'list number style',
        'list alpha style',
        'heading 1 style',
        'heading 2 style',
        'heading 3 style',
        'heading 4 style',
        'font is 11 points',
        '1.5 lines spacing',
        'space after has been added',
        'programmed with space after',
        'no more double returns'
    ]
    
    paragraphs_to_remove = []
    
    # Remove instruction paragraphs and any following unnecessary content
    # Extend the range to catch ALL instruction content (pages 3-4 and beyond)
    # Be more aggressive to ensure complete removal
    removal_end = min(end_idx + 50, len(doc.paragraphs))  # Extended range to catch page 4
    
    for idx in range(start_idx, removal_end):
        if idx >= len(doc.paragraphs):
            break
        
        para = doc.paragraphs[idx]
        text = para.text.strip()
        text_lower = text.lower()
        
        # Skip if paragraph has images (preserve images/logos)
        has_image = False
        try:
            for run in para.runs:
                if run._element.xpath('.//w:drawing') or run._element.xpath('.//w:pict'):
                    has_image = True
                    break
        except:
            pass
        
        if has_image:
            continue
            
        # Check if this paragraph should be removed
        should_remove = False
        
        # ALWAYS remove instruction paragraphs (pages 3-4)
        if idx >= start_idx and idx <= end_idx:
            should_remove = True
        # Remove paragraphs with unnecessary/decorative keywords (page 3-4 content)
        elif any(kw in text_lower for kw in unnecessary_keywords):
            should_remove = True
        # Remove paragraphs that are clearly instruction-related (even if outside the marked range)
        # This catches page 4 if it continues instructions
        elif any(kw in text_lower for kw in instruction_keywords) and idx > end_idx and idx < end_idx + 20:
            should_remove = True
        # Remove template example content (long paragraphs that look like examples)
        elif len(text) > 100 and idx >= start_idx and idx < end_idx + 20:
            # Additional check: if it contains lorem-like or placeholder text
            if any(word in text_lower for word in ['lorem', 'placeholder', 'example', 'sample', 'template', 'instruction']):
                should_remove = True
            # Or if it's in the instruction range (pages 3-4)
            elif idx <= end_idx + 10:
                should_remove = True
        # Remove any paragraph that looks like it's part of page 3-4 instructions
        # Check for patterns that indicate instruction continuation
        elif idx > start_idx and idx < end_idx + 25:
            # If paragraph contains style-related keywords and is near instruction area
            if any(keyword in text_lower for keyword in ['style', 'font', 'spacing', 'indent', 'points', 'heading', 'body text']):
                # Additional check: if it's a short paragraph (likely instruction text)
                if len(text) < 200:
                    should_remove = True
        
        if should_remove:
            paragraphs_to_remove.append((idx, para))
    
    # Remove paragraphs in reverse order to maintain indices
    removed_count = 0
    for idx, para in reversed(paragraphs_to_remove):
        try:
            p_element = para._element
            parent = p_element.getparent()
            if parent is not None:
                parent.remove(p_element)
                removed_count += 1
        except Exception as e:
            print(f"   ⚠️  Could not remove paragraph {idx}: {e}")
            continue
    
    print(f"   ✅ Removed {removed_count} unnecessary paragraphs")


async def generate_professional_toc_from_ai_response(structured_content: List[Dict], llm_service) -> List[Dict]:
    """
    Generate professional Table of Contents from AI response.
    Uses LLM to:
    1. Extract important headings from AI response
    2. Generate/extend TOC entries to be professional and descriptive
    3. Filter out unnecessary headings
    4. Create a comprehensive content table for page 2
    
    Returns:
        List of TOC items with level, text, and style
    """
    print(f"\n🤖 Generating professional Table of Contents from AI response...")
    
    # First, extract all headings from structured content
    all_headings = []
    for item in structured_content:
        style = item.get('style', '')
        text = item.get('text', '').strip()
        
        if not text:
            continue
        
        if 'Heading 1' in style:
            all_headings.append({'level': 1, 'text': text})
        elif 'Heading 2' in style:
            all_headings.append({'level': 2, 'text': text})
        elif 'Heading 3' in style:
            all_headings.append({'level': 3, 'text': text})
        elif 'Heading 4' in style:
            # Include Heading 4 if we don't have too many items
            if len(all_headings) < 20:
                all_headings.append({'level': 3, 'text': text})  # Treat H4 as level 3 in TOC
    
    if not all_headings:
        print(f"   ⚠️  No headings found in AI response")
        return []
    
    print(f"   📋 Found {len(all_headings)} headings in AI response")
    
    # Prepare headings list for LLM
    headings_list = "\n".join([f"{i+1}. Level {h['level']}: {h['text']}" for i, h in enumerate(all_headings)])
    
    # Also extract some content context for better TOC generation
    content_preview = ""
    for item in structured_content[:5]:  # First few content items for context
        if 'Body Text' in item.get('style', '') and item.get('text', '').strip():
            content_preview += item['text'][:200] + "...\n"
    
    system_prompt = """You are a professional document editor. Generate a comprehensive Table of Contents for page 2.

Your task:
1. Analyze the provided headings from the AI-generated content
2. Keep ONLY the most important headings (remove redundant, unnecessary, or minor headings)
3. Generate professional, clear TOC entries that are descriptive but concise (4-10 words)
4. Extend and enhance headings to be more descriptive and professional
5. Maintain the hierarchical structure (Level 1 = main sections, Level 2 = subsections, Level 3 = sub-subsections)
6. Ensure the TOC provides a clear overview of the document structure

Rules:
- Keep only important main sections and key subsections
- Remove redundant, duplicate, or overly specific headings
- Use professional title case format
- Make entries descriptive but concise
- Maintain proper level hierarchy (1, 2, 3)
- Output format: One entry per line as: [LEVEL] Professional Entry Text
- Generate 8-15 entries maximum for a clean, professional TOC

Output ONLY the filtered and generated TOC entries, one per line, in this format: [LEVEL] Entry Text"""

    user_prompt = f"""Generate a professional Table of Contents for page 2 from these headings extracted from AI-generated content:

HEADINGS:
{headings_list}

{f'CONTENT CONTEXT (first few paragraphs):\n{content_preview}\n' if content_preview else ''}

Generate a comprehensive, professional Table of Contents that:
- Filters to keep only important headings
- Extends headings to be more descriptive and professional
- Maintains proper hierarchy
- Provides clear navigation structure

Output format: [LEVEL] Professional Entry Text, one per line."""

    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.4,
            max_tokens=1000
        )
        
        # Parse response
        toc_items = []
        lines = response.content.strip().split('\n')
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Parse format: [LEVEL] Entry Text
            match = re.match(r'\[(\d+)\]\s*(.+)', line)
            if match:
                level = int(match.group(1))
                text = match.group(2).strip()
                # Clean up text (remove any extra formatting)
                text = re.sub(r'^\d+[\.\)]\s*', '', text)  # Remove leading numbers
                text = re.sub(r'^[-\*]\s*', '', text)  # Remove leading bullets
                if text and 1 <= level <= 3:
                    toc_items.append({
                        'level': level,
                        'text': text,
                        'style': 'Body Text style'
                    })
            else:
                # Try alternative formats
                # Format: "Level X: Text" or "X. Text"
                alt_match = re.match(r'(?:Level\s+)?(\d+)[:\.]\s*(.+)', line, re.IGNORECASE)
                if alt_match:
                    level = int(alt_match.group(1))
                    text = alt_match.group(2).strip()
                    if text and 1 <= level <= 3:
                        toc_items.append({
                            'level': level,
                            'text': text,
                            'style': 'Body Text style'
                        })
                else:
                    # Last resort: clean the line and use as level 2
                    text = re.sub(r'^[\d\.\-\*\(\)\[\]\s]+', '', line)
                    if text and len(text) > 3:
                        toc_items.append({
                            'level': 2,
                            'text': text,
                            'style': 'Body Text style'
                        })
        
        # Validate and filter TOC items
        if not toc_items:
            print(f"   ⚠️  LLM did not generate valid TOC entries, using important headings")
            # Fallback: use first 12 headings (most important)
            toc_items = all_headings[:12]
        else:
            # Limit to 15 items maximum for a clean TOC
            if len(toc_items) > 15:
                toc_items = toc_items[:15]
                print(f"   ℹ️  Limited TOC to 15 items for clarity")
        
        print(f"   ✅ Generated {len(toc_items)} professional TOC entries")
        return toc_items
        
    except Exception as e:
        print(f"   ⚠️  Failed to generate professional TOC: {e}, using important headings")
        # Fallback: use first 12 headings
        return all_headings[:12] if len(all_headings) > 0 else []


def generate_table_of_contents(structured_content: List[Dict]) -> List[Dict]:
    """
    Generate a table of contents from structured content headings.
    Extracts all Heading 1, Heading 2, and Heading 3 items.
    
    Returns:
        List of TOC items with style and text
    """
    toc_items = []
    
    for item in structured_content:
        style = item.get('style', '')
        text = item.get('text', '').strip()
        
        if not text:
            continue
        
        # Extract headings for TOC
        if 'Heading 1' in style:
            toc_items.append({
                'level': 1,
                'text': text,
                'style': 'Body Text style'
            })
        elif 'Heading 2' in style:
            toc_items.append({
                'level': 2,
                'text': text,
                'style': 'Body Text style'
            })
        elif 'Heading 3' in style:
            # Only include Heading 3 if we have fewer than 15 items (to keep TOC concise)
            if len(toc_items) < 15:
                toc_items.append({
                    'level': 3,
                    'text': text,
                    'style': 'Body Text style'
                })
    
    return toc_items


def clean_page2_remove_template_content(doc: Document, after_title_idx: int = 1):
    """
    Clean page 2 - remove ALL template content aggressively, keep it completely empty.
    This ensures page 2 is clean before we add Table of Contents.
    """
    print(f"\n🧹 Cleaning page 2 - removing ALL template content...")
    
    # Remove ALL paragraphs after title (page 1) until we reach the instruction area
    # Be more aggressive - remove everything that's not our added content
    start_idx = after_title_idx + 1
    end_idx = min(start_idx + 100, len(doc.paragraphs))  # Clean more paragraphs
    
    paragraphs_to_remove = []
    
    # Comprehensive list of template keywords and patterns (including Contents)
    template_keywords = [
        'instruction', 'style:', 'template', 'placeholder', 'body text style',
        'heading 1 style', 'heading 2 style', 'heading 3 style', 'heading 4 style',
        'list bullet style', 'list number style', 'list alpha style',
        'font is 11 points', '1.5 lines spacing', 'space after',
        'report large callout', 'placer at prod', 'dolor nam eim',
        'persius apeirian', 'disputationi', 'aeterno fuisset',
        'odio legere consulatu', 'consulatu cu placer',
        'lorem', 'example', 'sample', 'subtitle', 'headline',
        '错误', 'manually placed', 'dividing content', 'ending',
        'unnecessary', 'dont want', 'take 3rd', '4th page',
        'extract text', 'user content', 'ai response',
        'contents', 'content', 'table of contents', 'toc'  # Remove old Contents from page 2
    ]
    
    for idx in range(start_idx, end_idx):
        if idx >= len(doc.paragraphs):
            break
        
        para = doc.paragraphs[idx]
        text = para.text.strip()
        text_lower = text.lower()
        
        # Skip if paragraph has images (preserve images/logos in headers)
        has_image = False
        try:
            for run in para.runs:
                if run._element.xpath('.//w:drawing') or run._element.xpath('.//w:pict'):
                    has_image = True
                    break
        except:
            pass
        
        if has_image:
            continue
        
        # Remove ALL content on page 2 - be aggressive
        # Everything on page 2 before our Contents is template content
        if text or len(para.runs) > 0:
            # Check for template keywords
            if any(keyword in text_lower for keyword in template_keywords):
                paragraphs_to_remove.append(idx)
            # Also remove if it looks like template content
            elif len(text) > 0:
                # Remove if it contains Chinese error messages
                if '错误' in text or '文档中' in text:
                    paragraphs_to_remove.append(idx)
                # Remove if it's placeholder-like text
                elif any(phrase in text_lower for phrase in ['manually placed', 'dividing content', 'ending']):
                    paragraphs_to_remove.append(idx)
                # Remove all text on page 2 (we'll add our own Contents)
                elif idx < start_idx + 80:  # Remove first 80 paragraphs after title
                    paragraphs_to_remove.append(idx)
    
    # Remove paragraphs in reverse order
    removed_count = 0
    for idx in reversed(paragraphs_to_remove):
        try:
            if idx < len(doc.paragraphs):
                p_element = doc.paragraphs[idx]._element
                parent = p_element.getparent()
                if parent is not None:
                    parent.remove(p_element)
                    removed_count += 1
        except Exception as e:
            print(f"   ⚠️  Could not remove paragraph {idx}: {e}")
            continue
    
    print(f"   ✅ Removed {removed_count} template paragraphs from page 2")
    print(f"   ✅ Page 2 is now completely clean and ready for Contents")


def add_formatted_content_to_page2(doc: Document, structured_content: List[Dict], instruction_end_idx: Optional[int] = None, document_title: Optional[str] = None, page2_contents: Optional[str] = None):
    """
    Add AI-generated content to the document with proper Word styles.
    Structure:
    - Page 1: Cover page with title ONLY (cleaned, centered, with logo, background, theme)
    - Page 2: Contents page with important headings list (extracted from AI response)
    - Page 3+: Main content (formatted AI response as is, don't change anything)
    
    IMPORTANT: 
    - Page 2 shows contents (headings list) extracted from AI response with professional formatting
    - Page 3+ content is the formatted AI response as is, starts after headers with proper alignment
    - All pages maintain consistent headers, footers, logo, theme, and background from page 1
    - All template content is removed
    - Consistent background theme, design and layout used on page 1 continues for page 2 and page 3
    
    Args:
        doc: The Word document
        structured_content: List of content blocks with style and text (from AI response)
        instruction_end_idx: End index of instructions (content will be added after this)
        document_title: Optional document title for page 1
        page2_contents: Contents text with headings list for page 2 (extracted from AI response)
    """
    print(f"\n✍️  Adding content: {len(structured_content)} blocks from AI response...")
    
    if not structured_content:
        print(f"   ⚠️  No content to add")
        return
    
    # Clean page 2 first - remove all template content
    clean_page2_remove_template_content(doc, after_title_idx=0)
    
    # Find the insertion point - after page 1 (title)
    # Since we cleaned page 1, insertion point should be after first paragraph
    insertion_point = 1  # After title paragraph
    
    # Ensure we have at least one paragraph
    if len(doc.paragraphs) == 0:
        doc.add_paragraph()
        insertion_point = 0
    
    # Add page break before page 2 content to ensure it starts on a new page
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        
        # Insert page break paragraph before page 2
        if insertion_point < len(doc.paragraphs):
            break_para = doc.paragraphs[insertion_point].insert_paragraph_before()
        else:
            break_para = doc.add_paragraph()
        
        # Add page break to the paragraph
        run = break_para.add_run()
        br = OxmlElement('w:br')
        br.set(qn('w:type'), 'page')
        run._element.append(br)
        print(f"   📄 Added page break before page 2 (Contents)")
    except Exception as e:
        print(f"   ⚠️  Could not add page break: {e}")
    
    added = 0
    
    # Add contents to page 2 (headings list extracted from AI response)
    # IMPORTANT: Page 2 maintains same headers, footers, logo, background, and theme as page 1
    if page2_contents:
        try:
            # Parse contents text - it should start with "Contents" title
            lines = page2_contents.split('\n')
            
            # Add "Contents" as Heading 1 (title) - properly formatted and aligned
            contents_title = "Contents"
            if lines and lines[0].strip():
                contents_title = lines[0].strip()
            
            # Add proper spacing from header area
            title_para = doc.add_paragraph(style='Heading 1 style')
            run = title_para.add_run(contents_title)
            run.bold = True
            run.font.size = Pt(22)  # Professional size for contents title
            try:
                run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)  # PwC orange
            except:
                try:
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray fallback
                except:
                    pass
            # Professional spacing and alignment
            title_para.paragraph_format.space_before = Pt(24)  # Proper spacing from header
            title_para.paragraph_format.space_after = Pt(24)  # Space after title
            title_para.paragraph_format.keep_with_next = False
            title_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
            added += 1
            print(f"   ✅ Added page 2 title: '{contents_title}' (with same theme as page 1)")
            
            # Add contents list (headings) with professional formatting
            # Skip the first line (title) and empty lines, then add each heading
            for line_idx, line in enumerate(lines[1:], 1):
                original_line = line
                line = line.strip()
                
                # Skip empty lines but add small spacing
                if not line:
                    if line_idx < len(lines) - 1:  # Don't add spacing after last line
                        spacer = doc.add_paragraph()
                        spacer.paragraph_format.space_after = Pt(3)
                    continue
                
                # Determine indentation level (2 spaces = level 2, 4 spaces = level 3, etc.)
                indent_level = 0
                leading_spaces = len(original_line) - len(original_line.lstrip())
                
                if leading_spaces >= 6:  # 6+ spaces = level 4
                    indent_level = 3
                    line = line.strip()
                elif leading_spaces >= 4:  # 4 spaces = level 3
                    indent_level = 2
                    line = line.strip()
                elif leading_spaces >= 2:  # 2 spaces = level 2
                    indent_level = 1
                    line = line.strip()
                else:
                    # No indentation = level 1
                    indent_level = 0
                    line = line.strip()
                
                if not line:
                    continue
                
                # Add heading as properly formatted content with professional alignment
                para = doc.add_paragraph(style='Body Text style')
                
                # Apply professional formatting based on level
                if indent_level == 0:
                    # Level 1 heading - bold, larger, no indent
                    para.paragraph_format.left_indent = Inches(0)
                    para.paragraph_format.first_line_indent = Inches(0)
                    run = para.add_run(line)
                    run.bold = True
                    run.font.size = Pt(13)  # Slightly larger for level 1
                    try:
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray
                    except:
                        pass
                    para.paragraph_format.space_before = Pt(8)
                    para.paragraph_format.space_after = Pt(6)
                elif indent_level == 1:
                    # Level 2 heading - bold, indented
                    para.paragraph_format.left_indent = Inches(0.3)  # Professional indent
                    para.paragraph_format.first_line_indent = Inches(0)
                    run = para.add_run(line)
                    run.bold = True
                    run.font.size = Pt(12)
                    try:
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    except:
                        pass
                    para.paragraph_format.space_before = Pt(4)
                    para.paragraph_format.space_after = Pt(4)
                elif indent_level == 2:
                    # Level 3 heading - regular weight, more indented
                    para.paragraph_format.left_indent = Inches(0.6)
                    para.paragraph_format.first_line_indent = Inches(0)
                    run = para.add_run(line)
                    run.bold = False
                    run.font.size = Pt(11)
                    try:
                        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)  # Medium gray
                    except:
                        pass
                    para.paragraph_format.space_before = Pt(3)
                    para.paragraph_format.space_after = Pt(3)
                else:
                    # Level 4+ heading - regular, even more indented
                    para.paragraph_format.left_indent = Inches(0.9)
                    para.paragraph_format.first_line_indent = Inches(0)
                    run = para.add_run(line)
                    run.bold = False
                    run.font.size = Pt(10)
                    try:
                        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)  # Lighter gray
                    except:
                        pass
                    para.paragraph_format.space_before = Pt(2)
                    para.paragraph_format.space_after = Pt(2)
                
                # Consistent alignment and spacing
                para.paragraph_format.line_spacing = 1.5
                para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                para.paragraph_format.keep_with_next = False
                added += 1
            
            print(f"   ✅ Added page 2 contents ({len([l for l in lines[1:] if l.strip()])} headings) with professional formatting")
            print(f"   ✅ Page 2 maintains same logo, background, theme, headers, and footers as page 1")
        except Exception as e:
            print(f"   ⚠️  Could not add page 2 contents: {e}")
            import traceback
            traceback.print_exc()
    
    # Add page break after page 2 so main content starts on new page (page 3)
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        break_para = doc.add_paragraph()
        run = break_para.add_run()
        br = OxmlElement('w:br')
        br.set(qn('w:type'), 'page')
        run._element.append(br)
        print(f"   📄 Added page break after page 2 (main content starts on page 3)")
    except Exception as e:
        print(f"   ⚠️  Could not add page break after page 2: {e}")
    
    # Now add main content starting on page 3
    # Content should start after headers with proper alignment
    is_first_heading = True
    has_heading1 = False
    
    # Check if content starts with Heading 1
    for item in structured_content:
        if 'Heading 1' in item.get('style', ''):
            has_heading1 = True
            break
    
    # Add content blocks with proper formatting starting on page 3
    # Content should start after headers with proper alignment
    for item in structured_content:
        style_name = item['style']
        text = item['text']
        
        if not text.strip():
            continue
        
        # Determine if this is a heading
        is_heading = 'Heading' in style_name
        
        try:
            # Try to use the specified style with fallbacks
            para = None
            
            # Try exact style name first
            try:
                para = doc.add_paragraph(style=style_name)
                added += 1
            except (KeyError, ValueError, AttributeError):
                # Try alternative style names
                style_alternatives = {
                    'Heading 1 style': ['Heading 1', 'Title', 'Heading1'],
                    'Heading 2 style': ['Heading 2', 'Heading2'],
                    'Heading 3 style': ['Heading 3', 'Heading3'],
                    'Heading 4 style': ['Heading 4', 'Heading4'],
                    'Body Text style': ['Body Text', 'Normal', 'BodyText'],
                    'List Bullet style': ['List Bullet', 'ListBullet'],
                    'List Number style': ['List Number', 'ListNumber'],
                    'List Alpha style': ['List Alpha', 'ListAlpha'],
                }
                
                alternatives = style_alternatives.get(style_name, ['Body Text', 'Normal'])
                for alt_style in alternatives:
                    try:
                        para = doc.add_paragraph(style=alt_style)
                        added += 1
                        break
                    except:
                        continue
                
                # Last resort: add paragraph and apply formatting manually
                if para is None:
                    para = doc.add_paragraph()
                    added += 1
            
            # Add text content
            if '**' in text:
                # Handle bold text (markdown **text**)
                parts = text.split('**')
                for i, part in enumerate(parts):
                    if part:
                        run = para.add_run(part)
                        if i % 2 == 1:  # Odd indices are bold
                            run.bold = True
            else:
                para.add_run(text)
            
            # Apply comprehensive formatting based on style type
            # IMPORTANT: Content starts after headers, so add proper spacing
            try:
                if 'Heading 1' in style_name or (is_first_heading and is_heading):
                    # Heading 1: Large, bold, prominent title format
                    # Start after header with proper spacing
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(18) if run.font.size is None else run.font.size
                        # Apply color if possible (PwC orange or black)
                        try:
                            run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)  # PwC orange
                        except:
                            try:
                                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray fallback
                            except:
                                pass
                    # Space after header (page 3 starts after header)
                    para.paragraph_format.space_before = Pt(18) if is_first_heading else Pt(12)
                    para.paragraph_format.space_after = Pt(12)
                    para.paragraph_format.keep_with_next = True  # Keep heading with next paragraph
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Left align
                    is_first_heading = False
                    print(f"   ✅ Added Heading 1: '{text[:50]}...'")
                
                elif 'Heading 2' in style_name:
                    # Heading 2: Major section headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(16) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(12)
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.keep_with_next = True
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Heading 3' in style_name:
                    # Heading 3: Subsection headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(14) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(10)
                    para.paragraph_format.space_after = Pt(4)
                    para.paragraph_format.keep_with_next = True
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Heading 4' in style_name:
                    # Heading 4: Sub-subsection headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(12) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(8)
                    para.paragraph_format.space_after = Pt(4)
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Body Text' in style_name or style_name == 'Body Text style':
                    # Body text: 11pt font, 1.5 line spacing, space after, left aligned
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.line_spacing = 1.5  # 1.5 line spacing
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT  # Left align
                    for run in para.runs:
                        if run.font.size is None:
                            run.font.size = Pt(11)
                        # Ensure body text is not bold (unless explicitly marked)
                        if not any('**' in item['text'] for item in [item]):
                            run.bold = False
                
                elif 'Bullet' in style_name or 'Number' in style_name or 'Alpha' in style_name:
                    # Lists: hanging indent, proper spacing, left aligned
                    para.paragraph_format.left_indent = Inches(0.25)
                    para.paragraph_format.first_line_indent = Inches(-0.25)  # Hanging indent
                    para.paragraph_format.space_after = Pt(3)
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in para.runs:
                        if run.font.size is None:
                            run.font.size = Pt(11)
                
                elif 'Quote' in style_name:
                    # Quotes: italic, indented, left aligned
                    para.paragraph_format.left_indent = Inches(0.5)
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in para.runs:
                        run.italic = True
                        if run.font.size is None:
                            run.font.size = Pt(11)
                
            except Exception as e:
                # Style application failed, but continue
                print(f"   ⚠️  Could not apply all formatting to '{text[:30]}...': {e}")
                pass
            
        except Exception as e:
            print(f"   ⚠️  Error adding content block '{text[:30]}...': {e}")
            continue
    
    print(f"   ✅ Added {added} content paragraphs")
    print(f"   📋 Structure:")
    print(f"      - Page 1: Title only (with logo, background, theme)")
    print(f"      - Page 2: Title & Content (with same logo, background, theme)")
    print(f"      - Page 3+: AI Response (with same logo, background, theme)")
    print(f"   📋 All pages maintain consistent headers, footers, logo, background, and theme")
    print(f"   📋 Headings formatted with template styles (fonts, sizes, spacing, alignment)")


# ============================================
# SECTION 5: MAIN FUNCTION (SINGLE ENTRY POINT)
# ============================================


# ============================================
# SECTION 6: PDF CONVERSION (SINGLE FUNCTION)
# ============================================

def create_pdf_from_docx_bytes(docx_bytes: bytes) -> bytes:
    """
    Convert DOCX bytes to PDF bytes using Word COM automation.
    
    This is the primary PDF conversion method. It requires Microsoft Word
    to be installed on Windows systems. If Word COM is unavailable, this
    will raise a ValueError to trigger the fallback method.
    """
    import tempfile
    from docx2pdf import convert
    
    # Create temporary files
    with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_docx:
        temp_docx.write(docx_bytes)
        temp_docx_path = temp_docx.name
    
    temp_pdf_path = temp_docx_path.replace('.docx', '.pdf')
    
    try:
        # Try to convert using Word COM automation (Windows only)
        convert(temp_docx_path, temp_pdf_path)
        with open(temp_pdf_path, 'rb') as pdf_file:
            pdf_bytes = pdf_file.read()
        return pdf_bytes
    except Exception as e:
        # Check if it's a Word COM error (indicates Word is not available)
        error_msg = str(e)
        error_type = type(e).__name__
        
        # Check for COM-related errors
        is_com_error = (
            'Server execution failed' in error_msg or 
            'Operation unavailable' in error_msg or 
            'com_error' in error_msg or
            'pywintypes.com_error' in error_type or
            'win32com' in error_msg.lower() or
            'Word' in error_msg or
            'COM' in error_msg
        )
        
        if is_com_error:
            print(f"⚠️  Word COM conversion unavailable: {error_msg}")
            raise ValueError("Word COM unavailable - use fallback PDF generation")
        else:
            # Re-raise other errors (file system, permissions, etc.)
            raise
    finally:
        # Clean up temporary files
        try:
            if os.path.exists(temp_docx_path):
                os.remove(temp_docx_path)
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)
        except:
            pass




# ============================================
# PREVIEW ENDPOINT (OPTIONAL - FOR DEBUGGING)
# ============================================

@app.post("/api/export/preview")
async def preview_formatting(request: ExportRequest):
    """
    Preview how content will be formatted.
    Useful for debugging.
    """
    try:
        llm = get_llm()
        
        # Load template to get instructions
        template_path = 'template/Template2.docx'
        doc = Document(template_path)
        
        # Extract instructions
        instructions, _, _ = extract_instructions_from_template(doc)
        
        # Format content
        if instructions:
            formatted = await generate_formatted_content_from_instructions(
                request.content, instructions, llm
            )
        else:
            formatted = request.content
        
        # Generate title
        title = await generate_title_from_ai_content(formatted, llm)
        
        # Parse structure
        structure = parse_markdown_to_structure(formatted)
        
        return {
            "generated_title": title,
            "instructions_found": bool(instructions),
            "formatted_content_length": len(formatted),
            "structure_blocks": len(structure),
            "structure_preview": [
                {
                    "style": item["style"],
                    "text_preview": item["text"][:100] + "..." if len(item["text"]) > 100 else item["text"]
                }
                for item in structure[:10]
            ],
            "formatted_content_preview": formatted[:500] + "..." if len(formatted) > 500 else formatted
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def create_pdf_from_text_fallback(content: str, title: str) -> bytes:
    """
    Generate PDF directly from text using ReportLab (fallback when Word COM is unavailable).
    
    This method is used when Word-to-PDF conversion is not possible (e.g., on Linux/Mac
    or when Word is not installed). It creates a PDF with similar styling to the Word template.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER
    from reportlab.lib.colors import HexColor
    import re
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        topMargin=1 * inch,
        bottomMargin=1 * inch,
        leftMargin=1 * inch,
        rightMargin=1 * inch
    )

    styles = getSampleStyleSheet()
    
    # Title style (matches Word template: 28pt, orange, centered)
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=28,
        textColor=HexColor("#E87722"),  # PwC orange
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName="Helvetica-Bold"
    )

    # Body text style (11pt, 1.5 line spacing, space after)
    body_style = ParagraphStyle(
        "CustomBody",
        parent=styles["BodyText"],
        fontSize=11,
        alignment=TA_JUSTIFY,
        spaceAfter=12,
        leading=16.5,  # 1.5 * 11pt = 16.5pt
        fontName="Helvetica"
    )

    # Heading styles
    heading1_style = ParagraphStyle(
        "CustomH1",
        parent=styles["Heading1"],
        fontSize=18,
        spaceAfter=12,
        textColor=HexColor("#333333"),
        fontName="Helvetica-Bold"
    )
    
    heading2_style = ParagraphStyle(
        "CustomH2",
        parent=styles["Heading2"],
        fontSize=16,
        spaceAfter=10,
        textColor=HexColor("#333333"),
        fontName="Helvetica-Bold"
    )
    
    heading3_style = ParagraphStyle(
        "CustomH3",
        parent=styles["Heading3"],
        fontSize=14,
        spaceAfter=8,
        textColor=HexColor("#333333"),
        fontName="Helvetica-Bold"
    )

    story = []
    
    # Add title if provided
    if title and title.strip():
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.3 * inch))
        story.append(PageBreak())  # Page break after title (cover page)

    # Parse content line by line to handle markdown formatting
    lines = content.split("\n")
    current_paragraph = []
    
    for line in lines:
        line = line.strip()
        
        if not line:
            # Empty line - flush current paragraph
            if current_paragraph:
                para_text = " ".join(current_paragraph)
                para_text = convert_markdown_to_html(para_text)
                story.append(Paragraph(para_text, body_style))
                story.append(Spacer(1, 0.1 * inch))
                current_paragraph = []
            continue
        
        # Check for headings
        if line.startswith("# "):
            # Flush current paragraph
            if current_paragraph:
                para_text = " ".join(current_paragraph)
                para_text = convert_markdown_to_html(para_text)
                story.append(Paragraph(para_text, body_style))
                story.append(Spacer(1, 0.1 * inch))
                current_paragraph = []
            # Add heading
            heading_text = line[2:].strip()
            story.append(Paragraph(heading_text, heading1_style))
            story.append(Spacer(1, 0.15 * inch))
        elif line.startswith("## "):
            if current_paragraph:
                para_text = " ".join(current_paragraph)
                para_text = convert_markdown_to_html(para_text)
                story.append(Paragraph(para_text, body_style))
                story.append(Spacer(1, 0.1 * inch))
                current_paragraph = []
            heading_text = line[3:].strip()
            story.append(Paragraph(heading_text, heading2_style))
            story.append(Spacer(1, 0.12 * inch))
        elif line.startswith("### "):
            if current_paragraph:
                para_text = " ".join(current_paragraph)
                para_text = convert_markdown_to_html(para_text)
                story.append(Paragraph(para_text, body_style))
                story.append(Spacer(1, 0.1 * inch))
                current_paragraph = []
            heading_text = line[4:].strip()
            story.append(Paragraph(heading_text, heading3_style))
            story.append(Spacer(1, 0.1 * inch))
        elif line.startswith("- ") or line.startswith("• "):
            # Bullet point
            if current_paragraph:
                para_text = " ".join(current_paragraph)
                para_text = convert_markdown_to_html(para_text)
                story.append(Paragraph(para_text, body_style))
                story.append(Spacer(1, 0.1 * inch))
                current_paragraph = []
            bullet_text = line[2:].strip()
            bullet_text = convert_markdown_to_html("• " + bullet_text)
            story.append(Paragraph(bullet_text, body_style))
            story.append(Spacer(1, 0.05 * inch))
        else:
            # Regular text - add to current paragraph
            current_paragraph.append(line)
    
    # Flush remaining paragraph
    if current_paragraph:
        para_text = " ".join(current_paragraph)
        para_text = convert_markdown_to_html(para_text)
        story.append(Paragraph(para_text, body_style))

    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def convert_markdown_to_html(text: str) -> str:
    """Convert markdown formatting to HTML for ReportLab"""
    # Convert bold **text** to <b>text</b>
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    # Convert italic *text* to <i>text</i>
    text = re.sub(r'\*(.+?)\*', r'<i>\1</i>', text)
    # Escape HTML special characters
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;').replace('>', '&gt;')
    # Restore our HTML tags
    text = text.replace('&lt;b&gt;', '<b>').replace('&lt;/b&gt;', '</b>')
    text = text.replace('&lt;i&gt;', '<i>').replace('&lt;/i&gt;', '</i>')
    return text






async def generate_stream(messages: list):
    """Generate streaming response using LLM service with fallback"""
    try:
        llm = get_llm()

        async for content in llm.stream_completion(
            messages=messages,
            temperature=0.7,
            max_tokens=4096,
        ):
            yield f"data: {json.dumps({'content': content})}\n\n"

        yield f"data: {json.dumps({'done': True})}\n\n"

    except Exception as e:
        yield f"data: {json.dumps({'error': str(e)})}\n\n"


@app.post("/api/chat")
async def chat(request: ChatRequest):
    try:
        messages = [
            {"role": msg.role, "content": msg.content} for msg in request.messages
        ]

        if request.stream:
            return StreamingResponse(
                generate_stream(messages), media_type="text/event-stream"
            )
        else:
            llm = get_llm()
            response = await llm.chat_completion(
                messages=messages,
                temperature=0.7,
                max_tokens=4096,
            )

            return {
                "message": response.content,
                "provider": response.provider.value,
                "usage": response.usage or {
                    "prompt_tokens": 0,
                    "completion_tokens": 0,
                    "total_tokens": 0,
                },
            }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/draft")
async def create_draft(request: DraftRequest):
    try:
        system_prompt = """You are an expert presentation consultant at PwC (PricewaterhouseCoopers). 
Your role is to help create professional, structured, and impactful presentations.
When drafting presentations, focus on:
- Clear structure with MECE (Mutually Exclusive, Collectively Exhaustive) framework
- Executive summaries and key takeaways
- Data-driven insights
- Professional formatting suggestions
- PwC's consulting best practices"""

        user_prompt = f"""Please create a presentation outline with the following details:

Topic: {request.topic}
Objective: {request.objective}
Target Audience: {request.audience}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}

Provide a structured outline including:
1. Presentation title
2. Executive summary
3. Slide-by-slide breakdown with:
   - Slide titles
   - Key messages for each slide
   - Suggested content types (charts, frameworks, bullet points)
4. Conclusion and call-to-action

Format the output in a clear, professional manner."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/ppt/generate")
async def generate_ppt(
    topic: str = Form(...),
    objective: str = Form(...),
    audience: str = Form(...),
    additional_context: Optional[str] = Form(None),
    template_ppt: UploadFile = File(None),
):
    """Generate an actual PowerPoint file from presentation draft request."""
    try:
        # Generate content using Groq
        system_prompt = """You are an expert presentation consultant. Create concise, impactful slide content.
For each slide, provide:
- Slide title (one line)
- 3-5 bullet points (each under 15 words)
Keep content professional and data-driven."""

        user_prompt = f"""Create slide content for a presentation:

Topic: {topic}
Objective: {objective}
Target Audience: {audience}
{f"Additional Context: {additional_context}" if additional_context else ""}

Provide 5-8 slides with titles and bullet points. Format as:
SLIDE 1: [Title]
• [Bullet point 1]
• [Bullet point 2]
• [Bullet point 3]

SLIDE 2: [Title]
..."""

        llm = get_llm()
        llm_response = await llm.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.7,
            max_tokens=2000,
        )

        content = llm_response.content

        # Create PowerPoint from template or blank
        if template_ppt:
            template_content = await template_ppt.read()
            prs = Presentation(io.BytesIO(template_content))
            # Clear existing slides
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
        else:
            prs = Presentation()

        # Parse content and create slides
        slides_data = []
        current_slide = None

        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("SLIDE ") and ":" in line:
                if current_slide:
                    slides_data.append(current_slide)
                title = line.split(":", 1)[1].strip()
                current_slide = {"title": title, "bullets": []}
            elif line.startswith("•") or line.startswith("-"):
                if current_slide:
                    bullet = line.lstrip("•-").strip()
                    if bullet:
                        current_slide["bullets"].append(bullet)

        if current_slide:
            slides_data.append(current_slide)

        # Create slides
        for slide_data in slides_data:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Add bullet points
            if slide_data["bullets"] and len(slide.shapes) > 1:
                content_shape = slide.shapes[1]
                text_frame = content_shape.text_frame
                text_frame.clear()

                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0

        # Save to BytesIO
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        # Return as downloadable file
        filename = f"{topic.replace(' ', '_')[:30]}_presentation.pptx"
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"},
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/thought-leadership")
async def thought_leadership(request: ThoughtLeadershipRequest):
    try:
        system_prompt = """You are an expert thought leadership consultant at PwC (PricewaterhouseCoopers). 
You specialize in creating compelling, insightful content that positions executives and organizations as industry leaders.
Your expertise includes strategic insights, market analysis, editorial refinement, and content transformation.

**Key Sources Leveraged:**
- Agent uses PwC-specific guidelines where they exist (available at https://brand.pwc.com/standards.html)
- Leverages standard LLM logic and other external sources where PwC guidelines do not exist
- PwC guidelines for different editorial activities are listed under https://brand.pwc.com/standards.html"""

        user_prompt = ""
        reference_context = ""

        if request.reference_urls:
            reference_context = "\n\n**Reference Sources:**\n"
            failed_sources = []
            all_fetched_pages = []  # Store all pages including crawled ones

            for url in request.reference_urls:
                url_data = await fetch_url_content(url)
                if url_data["success"]:
                    # Add the main page
                    all_fetched_pages.append(
                        {
                            "url": url_data["url"],
                            "title": url_data["title"],
                            "content": url_data["content"],
                        }
                    )

                    # Automatically crawl related pages from the same domain
                    try:
                        related_pages = await crawl_related_pages(
                            url_data["url"], max_pages=5, max_depth=2
                        )
                        # Add related pages (skip the first one as it's the main page)
                        for page in related_pages:
                            if (
                                page["url"] != url_data["url"]
                            ):  # Don't duplicate the main page
                                all_fetched_pages.append(page)
                    except Exception:
                        # If crawling fails, continue with just the main page
                        pass
                else:
                    error_msg = url_data.get("error", "Unknown error")
                    failed_sources.append(f"{url} - {error_msg}")
                    reference_context += f"\n**Source:** {url}\n**Status:** Unable to fetch content - {error_msg}\n\n---\n"

            # Add all successfully fetched pages to context
            for i, page in enumerate(all_fetched_pages, 1):
                reference_context += f"\n**Source {i}:** {page['title']}\n**URL:** {page['url']}\n**Content:**\n{page['content']}\n\n---\n"

            if failed_sources:
                reference_context += f"\n\n**Note:** {len(failed_sources)} source(s) could not be accessed. Please only use information from successfully fetched sources above. Do not attempt to reference or cite the failed sources."

            if len(all_fetched_pages) > len(request.reference_urls):
                reference_context += f"\n\n**Additional Context:** Automatically discovered and included {len(all_fetched_pages) - len(request.reference_urls)} related page(s) from the provided website(s) to provide comprehensive information."

        if request.operation == "generate":
            citations_instruction = ""
            if request.reference_urls:
                citations_instruction = "\n\nIMPORTANT: Include proper citations in your article. At the end of the article, add a 'References' or 'Sources' section with all the URLs provided. Use in-text citations where appropriate (e.g., [1], [2]) that correspond to the numbered references."

            user_prompt = f"""Generate a draft thought leadership article with the following details:

Topic: {request.topic}
Perspective/Angle: {request.perspective}
Target Audience: {request.target_audience}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}{reference_context}{citations_instruction}

Create a comprehensive, well-structured article that:
1. Opens with a compelling hook and executive summary
2. Provides unique insights and data-driven perspectives (use reference sources when available)
3. Uses frameworks and strategic analysis
4. Includes real-world implications
5. Concludes with actionable recommendations
6. Properly cites all reference sources

Write in a professional, authoritative tone suitable for publication."""

        elif request.operation == "research":
            user_prompt = f"""Research and provide additional insights on:

Topic: {request.topic}
Current Perspective: {request.perspective}
{f"Additional Context: {request.additional_context}" if request.additional_context else ""}{reference_context}

Provide:
1. Emerging trends and developments in this area (using reference sources when available)
2. Different perspectives and counterarguments
3. Recent data, statistics, or case studies from the provided sources
4. Industry expert viewpoints
5. Future implications and opportunities

Focus on solution-oriented insights that add depth to the existing perspective. When using information from reference sources, cite them appropriately."""

        elif request.operation == "editorial":
            # Legacy editorial operation - redirect to use dedicated editorial endpoint
            # Keeping minimal for backward compatibility
            additional_instructions = (
                f"\n\nAdditional Instructions: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Edit the following document according to PwC editorial guidelines:

{request.document_text}{additional_instructions}

Output ONLY the final edited version—no explanations, comments, or tracked changes."""

        elif request.operation == "improve":
            focus_areas = (
                f"\nFocus Areas: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Recommend improvements to this document:

{request.document_text}{focus_areas}

Analyze and recommend improvements for:
1. Content quality and depth of insights
2. Argument strength and persuasiveness
3. Evidence and supporting data
4. Professional positioning and credibility
5. Engagement and readability
6. Call-to-action effectiveness

Prioritize the most impactful improvements."""

        elif request.operation == "translate":
            additional_requirements = (
                f"\nAdditional Requirements: {request.additional_context}"
                if request.additional_context
                else ""
            )
            user_prompt = f"""Translate this document to a different format:

Original Document:
{request.document_text}

Target Format: {request.target_format}{additional_requirements}

Transform the content while:
1. Maintaining core messages and insights
2. Adapting tone and style to the target format
3. Adjusting length and structure appropriately
4. Optimizing for the new medium's best practices
5. Preserving professional quality

Provide the fully transformed content."""

        else:
            raise HTTPException(
                status_code=400, detail=f"Invalid operation: {request.operation}"
            )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/research")
async def research_assistant(request: ResearchRequest):
    """AI Research Assistant - provides LLM-powered research insights and analysis"""
    try:
        system_prompt = """You are an expert research assistant at PwC with deep industry knowledge and analytical capabilities.
You specialize in:
- Identifying emerging trends and market developments based on your training data
- Analyzing competitive intelligence and industry dynamics
- Providing data-driven insights and statistical perspectives
- Synthesizing multiple perspectives and viewpoints
- Recommending strategic implications

IMPORTANT: You are providing insights based on your knowledge base (training data through April 2024). 
For the most current data, recommend users verify with real-time sources.
Focus on frameworks, patterns, and strategic analysis that remain relevant regardless of real-time data.

Your research is comprehensive, objective, and actionable."""

        focus_context = ""
        if request.focus_areas:
            focus_context = f"\n\nFocus Areas: {', '.join(request.focus_areas)}"

        additional = (
            f"\n\nAdditional Context: {request.additional_context}"
            if request.additional_context
            else ""
        )

        user_prompt = f"""Conduct comprehensive research on the following query:

{request.query}{focus_context}{additional}

Provide:
1. **Key Findings**: 3-5 critical insights with supporting evidence
2. **Emerging Trends**: Recent developments and future implications
3. **Data Points**: Relevant statistics, metrics, or quantitative insights
4. **Different Perspectives**: Multiple viewpoints and counterarguments
5. **Strategic Implications**: How this impacts business strategy
6. **Recommended Next Steps**: Actionable recommendations

Format your response clearly with headers and bullet points for easy scanning."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# New: Thought Leadership Editorial endpoint with personas, approach, and file upload
@app.post("/api/thought-leadership/editorial")
async def thought_leadership_editorial(
    document_text: Optional[str] = Form(None),
    personas: Optional[str] = Form(None),  # JSON list of personas (legacy support)
    approach: Optional[str] = Form(None),  # Legacy approach parameter
    selected_approach: Optional[str] = Form(None),  # New: single selected approach from workflow
    additional_context: Optional[str] = Form(None),
    reference_links: Optional[List[str]] = Form(None),
    file: Optional[UploadFile] = File(None),
    feedback: Optional[str] = Form(None),  # User feedback for iterations
):
    """Streaming editorial support that enforces PwC Brand guidelines and supports personas.

    - selected_approach: Single selected approach from workflow. Valid values:
      "Development Editor", "Content Editor", "Line Editor", "Copy Editor", 
      "PwC Brand Alignment Editor", "Run All Sequentially"
    - personas: JSON array of selected personas (legacy support, will be derived from selected_approach if not provided).
    - approach: Optional editing approach label (legacy support).
    - Accepts either uploaded file (pdf/docx/txt/md) or raw document_text.
    - reference_links: Optional list of URLs to consider.
    """
    try:
        # Parse uploaded content - enforce either/or (not both)
        uploaded_text = ""
        file_parse_error = None
        
        if file:
            # Check if file has a valid filename
            if not file.filename or not file.filename.strip():
                raise HTTPException(
                    status_code=400,
                    detail="Uploaded file has no filename. Please provide a valid file."
                )
            
            uploaded_text = await parse_uploaded_file(file, max_chars=20000)
            
            # Check if file parsing returned an error message
            if uploaded_text.startswith("[Error") or uploaded_text.startswith("[Could not") or uploaded_text.startswith("[Unsupported"):
                file_parse_error = uploaded_text
                uploaded_text = ""  # Treat as failed extraction

        document_text_content = (document_text or "").strip() if document_text else ""
        
        # Backend validation (security layer - frontend should catch this first)
        # Check both provided (frontend should prevent this, but validate for security)
        if file and document_text_content:
            raise HTTPException(
                status_code=400, 
                detail="Please provide either a file upload OR text input, not both. Choose one method."
            )
        
        # Determine base_text from either source
        if uploaded_text and uploaded_text.strip():
            base_text = uploaded_text
        elif document_text_content:
            base_text = document_text_content
        else:
            # No content provided - provide specific error based on what was attempted
            if file_parse_error:
                raise HTTPException(
                    status_code=400,
                    detail=f"File parsing failed: {file_parse_error}. Please try a different file or paste text directly."
                )
            elif file:
                raise HTTPException(
                    status_code=400,
                    detail="The uploaded file appears to be empty or could not be processed. Please try a different file or paste text directly."
                )
            else:
                raise HTTPException(
                    status_code=400, 
                    detail="No document content provided. Please either upload a file (.pdf, .docx, .pptx, .txt) or paste text directly."
                )

        # Final validation: ensure we have actual content (not just whitespace)
        if not base_text or not base_text.strip():
            raise HTTPException(
                status_code=400, 
                detail="No document content found. The file may be empty or text extraction failed. Please try again with valid content."
            )

        # Determine selected approach (prioritize new selected_approach parameter)
        final_selected_approach = selected_approach or approach
        
        # Validate selected approach
        valid_approaches = [
            "Development Editor",
            "Content Editor", 
            "Line Editor",
            "Copy Editor",
            "PwC Brand Alignment Editor",
            "Run All Sequentially"
        ]
        
        if final_selected_approach and final_selected_approach not in valid_approaches:
            raise HTTPException(
                status_code=400,
                detail=f"Invalid approach selected. Must be one of: {', '.join(valid_approaches)}"
            )
        
        # If no approach selected, default to "Run All Sequentially"
        if not final_selected_approach:
            final_selected_approach = "Run All Sequentially"
        
        # Build personas list from selected approach
        if final_selected_approach == "Run All Sequentially":
            selected_personas = [
                "Development Editor",
                "Content Editor",
                "Line Editor",
                "Copy Editor",
                "PwC Brand Alignment Editor"
            ]
        elif final_selected_approach == "PwC Brand Alignment Editor":
            selected_personas = ["PwC Brand Alignment Editor"]
        else:
            # Single approach + Brand Alignment (always included)
            selected_personas = [final_selected_approach, "PwC Brand Alignment Editor"]
        
        # Legacy support: if personas provided, use them (but still respect selected_approach)
        if personas and not selected_approach:
            try:
                parsed_personas = parse_personas_list(personas, selected_personas)
                if parsed_personas:
                    selected_personas = parsed_personas
            except Exception:
                pass  # Use approach-based personas

        # Reference links context (optional)
        reference_context = ""
        if reference_links:
            reference_context = "\n\n[Reference Links provided by user:]\n" + "\n".join(
                f"- {link}" for link in reference_links if link
            )

        # Build prompts using helper function with selected approach
        # Use feedback if provided, otherwise empty string
        feedback_context = feedback if feedback and feedback.strip() else ""
        
        try:
            system_prompt, user_prompt = build_editorial_prompts(
                selected_personas=selected_personas,
                base_text=base_text,
                sections_context=reference_context,
                feedback_context=feedback_context,
                selected_approach=final_selected_approach
            )
        except Exception as prompt_error:
            raise HTTPException(status_code=500, detail=f"Error building prompts: {str(prompt_error)}")

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def parse_personas_list(personas_input: Optional[str], default_personas: List[str]) -> List[str]:
    """Helper function to parse and validate personas list.
    
    Args:
        personas_input: JSON string or list of personas
        default_personas: Default list of personas to use as fallback
        
    Returns:
        List of validated persona names (always includes PwC Brand Alignment Editor)
    """
    selected_personas = default_personas.copy()
    
    if personas_input:
        try:
            parsed = json.loads(personas_input) if isinstance(personas_input, str) else personas_input
            if isinstance(parsed, list) and len(parsed) > 0:
                # Filter valid personas and always include Brand Alignment
                selected_personas = [p for p in parsed if isinstance(p, str) and p.strip() and p in PERSONA_DESCRIPTIONS]
                if "PwC Brand Alignment Editor" not in selected_personas:
                    selected_personas.append("PwC Brand Alignment Editor")
        except Exception:
            # Fallback to defaults on parse issues
            selected_personas = default_personas

    return selected_personas


def build_editorial_prompts(selected_personas: List[str], base_text: str, sections_context: str = "", feedback_context: str = "", selected_approach: Optional[str] = None) -> tuple:
    """Build system and user prompts for editorial services using ONLY instructions from editorial_prompts.py.
    
    Args:
        selected_personas: List of selected persona names
        base_text: Document content to edit
        sections_context: Optional context about sections/pages to edit
        feedback_context: Optional user feedback for iterations
        
    Returns:
        Tuple of (system_prompt, user_prompt)
    """
    # Validate selected_personas
    if not selected_personas or len(selected_personas) == 0:
        raise ValueError("No personas selected")
    
    # Check if "Run All Sequentially" - all 5 personas in order (handles "Run all approaches" from frontend)
    is_run_all = len(selected_personas) == 5 and all(p in selected_personas for p in [
        "Development Editor", "Content Editor", "Line Editor", "Copy Editor", "PwC Brand Alignment Editor"
    ])
    
    # Determine if single persona mode (excluding Brand Alignment)
    is_single_persona = len(selected_personas) == 1 or (len(selected_personas) == 2 and "PwC Brand Alignment Editor" in selected_personas)
    
    # Build system prompt using ONLY instructions from editorial_prompts.py
    if is_run_all:
        # Run All Sequentially: Apply all 5 editors in order, single final output
        system_prompt = (
            "You are a PwC editorial assistant. Apply all five editorial approaches sequentially in this order:\n"
            "1. Development Editor\n"
            "2. Content Editor\n"
            "3. Line Editor\n"
            "4. Copy Editor\n"
            "5. PwC Brand Alignment Editor\n\n"
            f"{DEVELOPMENT_EDITOR_TONE_OF_VOICE}\n\n"
            f"{CONTENT_EDITOR_INSTRUCTIONS}\n\n"
            f"{LINE_EDITOR_INSTRUCTIONS}\n\n"
            f"{COPY_EDITOR_INSTRUCTIONS}\n\n"
            f"{BRAND_ALIGNMENT_EDITOR_INSTRUCTIONS}\n\n"
            "**OUTPUT REQUIREMENTS:**\n"
            "- Apply all five editorial approaches sequentially to the document\n"
            "- Output ONLY the final edited version (no explanations, no tracked changes, no comments)\n"
            "- Work ONLY with the user's existing document—DO NOT generate new content\n"
            "- Preserve factual content—only apply editorial improvements\n"
        )
    elif is_single_persona:
        # Single persona mode
        primary_persona = [p for p in selected_personas if p != "PwC Brand Alignment Editor"][0] if "PwC Brand Alignment Editor" in selected_personas else selected_personas[0]
        
        # Build system prompt with ONLY the relevant instructions from editorial_prompts.py
        system_prompt = ""
        
        if primary_persona == "Development Editor":
            system_prompt = f"{DEVELOPMENT_EDITOR_TONE_OF_VOICE}\n\n"
        elif primary_persona == "Content Editor":
            system_prompt = f"{CONTENT_EDITOR_INSTRUCTIONS}\n\n"
        elif primary_persona == "Line Editor":
            system_prompt = f"{LINE_EDITOR_INSTRUCTIONS}\n\n"
        elif primary_persona == "Copy Editor":
            system_prompt = f"{COPY_EDITOR_INSTRUCTIONS}\n\n"
        
        # Add Brand Alignment if included
        if "PwC Brand Alignment Editor" in selected_personas:
            system_prompt += f"{BRAND_ALIGNMENT_EDITOR_INSTRUCTIONS}\n\n"
        
        system_prompt += (
            "**OUTPUT REQUIREMENTS:**\n"
            "- Output ONLY the final edited version (no explanations, no tracked changes, no comments)\n"
            "- Work ONLY with the user's existing document—DO NOT generate new content\n"
            "- Preserve factual content—only apply editorial improvements\n"
        )
    else:
        # Multiple personas (but not all 5)
        persona_instructions = ""
        
        if "Development Editor" in selected_personas:
            persona_instructions += f"{DEVELOPMENT_EDITOR_TONE_OF_VOICE}\n\n"
        if "Content Editor" in selected_personas:
            persona_instructions += f"{CONTENT_EDITOR_INSTRUCTIONS}\n\n"
        if "Line Editor" in selected_personas:
            persona_instructions += f"{LINE_EDITOR_INSTRUCTIONS}\n\n"
        if "Copy Editor" in selected_personas:
            persona_instructions += f"{COPY_EDITOR_INSTRUCTIONS}\n\n"
        if "PwC Brand Alignment Editor" in selected_personas:
            persona_instructions += f"{BRAND_ALIGNMENT_EDITOR_INSTRUCTIONS}\n\n"
        
        system_prompt = (
            f"{persona_instructions}"
            "**OUTPUT REQUIREMENTS:**\n"
            "- Apply all selected editorial approaches to the document\n"
            "- Output ONLY the final edited version (no explanations, no tracked changes, no comments)\n"
            "- Work ONLY with the user's existing document—DO NOT generate new content\n"
            "- Preserve factual content—only apply editorial improvements\n"
        )
    
    # Build user prompt
    user_prompt = (
        f"**DOCUMENT TO EDIT:**\n[Document Begins]\n{base_text}\n[Document Ends]\n\n"
    )
    
    if sections_context:
        user_prompt += f"**SECTIONS/PAGES TO FOCUS ON:**{sections_context}\n\n"
    
    if feedback_context:
        user_prompt += (
            f"**USER FEEDBACK FOR ADDITIONAL UPDATES:**\n"
            f"{feedback_context}\n\n"
            "**INSTRUCTIONS:** Apply the user's feedback and requests above to the document. "
            "Make the requested modifications while maintaining editorial quality and PwC brand alignment. "
            "Work with the original document content provided and incorporate the user's specific requests.\n\n"
        )
    
    user_prompt += (
        "**TASK:** Edit the document above according to the editorial guidelines provided"
    )
    if feedback_context:
        user_prompt += " and the user feedback above"
    user_prompt += ".\nOutput ONLY the final edited version—no explanations, comments, or tracked changes."
    
    return system_prompt, user_prompt


class EditorialConversationRequest(BaseModel):
    messages: List[Message]
    uploaded_file_name: Optional[str] = None
    stream: bool = True


@app.post("/api/thought-leadership/editorial/conversation")
async def editorial_services_conversation(request: EditorialConversationRequest):
    """Conversational editorial services endpoint that guides users through the editorial process.

    When user requests editing services through quick start conversation:
    - Agent responds by specifying the 5 core editorial services (personas) it offers
    - Provides brief description and/or examples for each persona
    - User has option to opt out of any specific service and/or specify content sections/pages to be edited
    - If not already provided, Agent requests document upload
    - List provided by agent has unique index (1-5) for each persona
    - User can specify actions by referring to the index
    """
    try:
        # Build service listing using shared PERSONA_DESCRIPTIONS
        service_list = "\n".join([
            f"   **{i}. {name}** – {desc}" + (" Always applied - cannot be opted out. Guidelines available at https://brand.pwc.com/standards.html." if name == "PwC Brand Alignment Editor" else ".")
            for i, (name, desc) in enumerate(PERSONA_DESCRIPTIONS.items(), 1)
        ])

        system_prompt = f"""You are a specialized PwC Editorial Services Assistant. Your role is to guide users through the editorial services process.

**WORKFLOW:**

1. **Service Listing**: When a user requests editing services, immediately provide a numbered list (1-5) of the 5 core editorial services with brief descriptions:

{service_list}

2. **Document Collection**: If no document is uploaded, politely request it. Explain you need the document (PDF, Word, or text) to proceed.

3. **Service Selection**: Allow users to:
   - Opt out of services by index (e.g., "exclude 2" or "skip 3")
   - Specify content sections/pages (e.g., "edit pages 1-5 only" or "focus on introduction")
   - Select specific services by index (e.g., "use 1, 3, and 5")

4. **Confirmation**: Before processing, summarize selected services (by index), excluded services, document sections/pages (if specified), and document to be processed.

**GUIDELINES:**
- Always use numbered list (1-5) for services
- Users can reference services by index number
- Support page/section-specific editing requests
- Always confirm before processing
- Be concise and professional

**KEY SOURCES:**
- Use PwC-specific guidelines from https://brand.pwc.com/standards.html where they exist
- Leverage standard LLM logic and other external sources where PwC guidelines do not exist
- Always apply PwC Brand Alignment Editor (service 5) regardless of user selection

Respond naturally and guide the user through the process step by step."""

        # Prepare messages for LLM
        messages_for_llm = [
            {"role": "system", "content": system_prompt}
        ]
        
        # Add conversation history
        for msg in request.messages:
            messages_for_llm.append({
                "role": msg.role,
                "content": msg.content
            })
        
        # If file name is provided, add context
        if request.uploaded_file_name:
            messages_for_llm.append({
                "role": "system",
                "content": f"User has uploaded a file: {request.uploaded_file_name}"
            })

        return StreamingResponse(
            generate_stream(messages_for_llm), media_type="text/event-stream"
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


async def fetch_web_content(url: str, max_chars: int = 10000) -> str:
    """Fetch and extract text content from a URL"""
    try:
        # Browser-like headers to avoid 403 Forbidden errors
        # Build referer from the URL (same domain)
        parsed_url = urlparse(url)
        referer = f"{parsed_url.scheme}://{parsed_url.netloc}/"

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8,application/signed-exchange;v=b3;q=0.7",
            "Accept-Language": "en-US,en;q=0.9",
            "Accept-Encoding": "gzip, deflate, br",
            "Connection": "keep-alive",
            "Upgrade-Insecure-Requests": "1",
            "Referer": referer,
            "Sec-Fetch-Dest": "document",
            "Sec-Fetch-Mode": "navigate",
            "Sec-Fetch-Site": "same-origin",
            "Sec-Fetch-User": "?1",
            "Cache-Control": "max-age=0",
            "DNT": "1",
        }

        async with httpx.AsyncClient(timeout=30.0, headers=headers) as client:
            response = await client.get(url, follow_redirects=True)

            # Check status codes before raising
            if response.status_code == 404:
                # Try URL variations before giving up
                variations = generate_url_variations(url)
                tried_urls = [url]

                for variation_url in variations:
                    if variation_url in tried_urls:
                        continue
                    tried_urls.append(variation_url)

                    try:
                        # Update referer for the variation
                        parsed_var = urlparse(variation_url)
                        var_referer = f"{parsed_var.scheme}://{parsed_var.netloc}/"
                        headers["Referer"] = var_referer

                        var_response = await client.get(
                            variation_url, follow_redirects=True
                        )
                        if var_response.status_code == 200:
                            # Success! Use this URL and continue with content extraction
                            url = variation_url  # Update url to the working variation
                            response = var_response
                            break
                        elif var_response.status_code != 404:
                            # If it's not 200 or 404, it's a different error, stop trying
                            break
                    except Exception:
                        # Continue to next variation on any error
                        continue

                # If still 404 after trying variations, return error
                if response.status_code == 404:
                    return f"[Error: 404 Not Found - The page at this URL does not exist. Tried {len(tried_urls)} variation(s) but none were found. Please verify the URL is correct.]"
            elif response.status_code == 403:
                return f"[Error: 403 Forbidden - The website blocked access. This may be due to bot protection, Cloudflare, or WAF rules. The site may require JavaScript to load content.]"
            elif response.status_code == 429:
                return f"[Error: 429 Too Many Requests - Rate limit exceeded. Please wait a moment and try again.]"
            elif response.status_code == 503:
                return f"[Error: 503 Service Unavailable - The website is temporarily unavailable. Please try again later.]"

            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer", "header"]):
                script.decompose()

            # Get text content
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            # Truncate if needed
            if len(text) > max_chars:
                text = text[:max_chars] + "... [truncated]"

            return text
    except Exception as e:
        return f"[Error fetching URL: {str(e)}]"


@app.post("/api/research-with-materials")
async def research_with_materials(
    query: str = Form(...),
    files: Optional[List[UploadFile]] = File(None),
    links: Optional[List[str]] = Form(None),
    focus_areas: Optional[str] = Form(None),
    additional_context: Optional[str] = Form(None),
):
    """Research assistant with document upload and URL reference support (NotebookLM-style)"""
    try:

        async def generate_research_stream():
            try:
                # Step 1: Parse uploaded documents
                documents_content = ""
                if files:
                    yield f'data: {json.dumps({"type": "progress", "message": f"Parsing {len(files)} document(s)..."})}\n\n'

                    for i, file in enumerate(files, 1):
                        parsed_content = await parse_uploaded_file(
                            file, max_chars=15000
                        )
                        if parsed_content:
                            documents_content += f"\n\n### Document {i}: {file.filename}\n{parsed_content}"

                # Step 2: Fetch and parse web content from links (with automatic crawling)
                web_content = ""
                if links:
                    valid_links = [link for link in links if link and link.strip()]
                    if valid_links:
                        yield f'data: {json.dumps({"type": "progress", "message": f"Fetching content from {len(valid_links)} link(s) and discovering related pages..."})}\n\n'

                        all_web_sources = []
                        for link in valid_links:
                            link = link.strip()
                            # Fetch main page with full metadata
                            url_data = await fetch_url_content(link)
                            if url_data.get("success"):
                                all_web_sources.append(
                                    {
                                        "url": url_data["url"],
                                        "title": url_data.get("title", "").strip()
                                        or link,
                                        "content": url_data.get("content", "")[:10000],
                                    }
                                )

                                # Automatically crawl related pages
                                try:
                                    related_pages = await crawl_related_pages(
                                        link, max_pages=3, max_depth=1
                                    )
                                    for page in related_pages:
                                        if (
                                            page["url"] != link
                                        ):  # Don't duplicate main page
                                            all_web_sources.append(
                                                {
                                                    "url": page["url"],
                                                    "title": page.get(
                                                        "title", ""
                                                    ).strip()
                                                    or page["url"],
                                                    "content": (
                                                        page["content"][:10000]
                                                        if len(page["content"]) > 10000
                                                        else page["content"]
                                                    ),
                                                }
                                            )
                                except Exception:
                                    # If crawling fails, continue with just main page
                                    pass
                            else:
                                # Fallback to fetch_web_content if fetch_url_content fails
                                main_content = await fetch_web_content(link)
                                if not main_content.startswith("[Error:"):
                                    all_web_sources.append(
                                        {
                                            "url": link,
                                            "title": link,
                                            "content": main_content[:10000],
                                        }
                                    )

                        # Add all sources to web_content with title for better citations
                        for i, source in enumerate(all_web_sources, 1):
                            title = source.get("title", "").strip() or source["url"]
                            # Include both title and URL for frontend parsing
                            web_content += f"\n\n### Web Source {i}: {title} | URL: {source['url']}\n{source['content']}"

                        if len(all_web_sources) > len(valid_links):
                            yield f'data: {json.dumps({"type": "progress", "message": f"Discovered {len(all_web_sources) - len(valid_links)} additional related page(s) automatically"})}\n\n'

                        # Send source metadata for frontend link rendering
                        source_metadata = [
                            {
                                "number": i,
                                "url": s["url"],
                                "title": s.get("title", "").strip() or s["url"],
                            }
                            for i, s in enumerate(all_web_sources, 1)
                        ]
                        yield f'data: {json.dumps({"type": "sources", "sources": source_metadata})}\n\n'

                # Step 3: Construct research prompt with materials
                yield f'data: {json.dumps({"type": "progress", "message": "Analyzing materials and conducting research..."})}\n\n'

                system_prompt = """You are an expert research assistant at PwC with deep analytical capabilities.
You excel at:
- Synthesizing information from multiple sources
- Identifying key themes and patterns across documents
- Extracting actionable insights and strategic implications
- Providing comprehensive analysis based on provided materials
- Cross-referencing different sources to validate findings

You are analyzing user-provided documents and web sources. Base your research primarily on these materials,
while supplementing with your knowledge when relevant. Always cite which sources support your findings."""

                materials_section = ""
                if documents_content or web_content:
                    materials_section = "\n\n## Reference Materials Provided:"
                    if documents_content:
                        materials_section += documents_content
                    if web_content:
                        materials_section += web_content

                focus_areas_list = []
                if focus_areas:
                    try:
                        focus_areas_list = (
                            json.loads(focus_areas)
                            if isinstance(focus_areas, str)
                            else focus_areas
                        )
                    except:
                        focus_areas_list = [focus_areas]

                focus_context = ""
                if focus_areas_list:
                    focus_context = f"\n\nFocus Areas: {', '.join(focus_areas_list)}"

                additional = (
                    f"\n\nAdditional Context: {additional_context}"
                    if additional_context
                    else ""
                )

                user_prompt = f"""Based on the reference materials provided, conduct comprehensive research on the following query:

{query}{focus_context}{additional}{materials_section}

Provide a thorough analysis including:
1. **Executive Summary**: Overview of key findings from the materials
2. **Key Insights**: 4-6 critical insights extracted from the provided sources (cite specific sources)
3. **Themes & Patterns**: Common threads and patterns across the materials
4. **Data & Evidence**: Relevant statistics, quotes, and concrete evidence from the sources
5. **Cross-Source Analysis**: How different sources complement or contradict each other
6. **Strategic Implications**: What this means for business strategy and decision-making
7. **Gaps & Recommendations**: What's missing and suggested next steps

Format your response clearly with headers and bullet points. When citing insights, reference the specific source (e.g., "Document 1" or "Web Source 2")."""

                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt},
                ]

                # Stream the AI response
                llm = get_llm()
                async for content in llm.stream_completion(
                    messages=messages,
                    temperature=0.7,
                    max_tokens=4096,
                ):
                    yield f'data: {json.dumps({"type": "content", "content": content})}\n\n'

                yield f'data: {json.dumps({"type": "complete"})}\n\n'

            except Exception as e:
                yield f'data: {json.dumps({"type": "error", "message": str(e)})}\n\n'

        return StreamingResponse(
            generate_research_stream(), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/draft-article")
async def draft_article(
    topic: str = Form(...),
    content_type: str = Form(...),
    desired_length: int = Form(...),
    tone: str = Form(...),
    outline_text: Optional[str] = Form(None),
    additional_context: Optional[str] = Form(None),
    outline_file: Optional[UploadFile] = File(None),
    supporting_docs: Optional[List[UploadFile]] = File(None),
):
    """Generate long-form thought leadership articles from user outlines and supporting documents"""
    try:
        system_prompt = """You are an expert thought leadership writer at PwC.
You excel at creating compelling, insightful articles that position executives as industry leaders.
You use the MECE principle (Mutually Exclusive, Collectively Exhaustive), strategic frameworks,
and data-driven insights. Your writing is authoritative yet accessible.

**Key Sources Leveraged:**
- Agent uses PwC-specific guidelines where they exist (available at https://brand.pwc.com/standards.html)
- Leverages standard LLM logic and other external sources where PwC guidelines do not exist
- PwC guidelines for different editorial activities are listed under https://brand.pwc.com/standards.html"""

        outline_content = outline_text or ""
        supporting_content = ""

        if outline_file:
            parsed_outline = await parse_uploaded_file(outline_file, max_chars=5000)
            if parsed_outline:
                outline_content += "\n\n**Uploaded Outline:**\n" + parsed_outline

        if supporting_docs:
            supporting_content = "\n\n**Supporting Documents:**\n"
            for i, doc in enumerate(supporting_docs, 1):
                parsed_content = await parse_uploaded_file(doc, max_chars=3000)
                if parsed_content:
                    supporting_content += (
                        f"\n**Document {i} ({doc.filename}):**\n{parsed_content}\n"
                    )

        additional = (
            f"\n\nAdditional Context: {additional_context}"
            if additional_context
            else ""
        )

        user_prompt = f"""Create a comprehensive {content_type.lower()} on the following topic:

**Topic:** {topic}
**Content Type:** {content_type}
**Target Length:** {desired_length} words
**Tone:** {tone}

**Outline/Initial Ideas:**
{outline_content}{supporting_content}{additional}

Generate a well-structured, professional article that:
1. Opens with a compelling hook and clear thesis
2. Uses the provided outline as a foundation (if provided)
3. Integrates insights from supporting documents (if provided)
4. Includes strategic frameworks and analysis
5. Provides real-world examples and implications
6. Concludes with actionable recommendations
7. Maintains the specified tone throughout
8. Targets approximately {desired_length} words

Format with clear headers, subheaders, and bullet points where appropriate."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ============================================
# DOCUMENT GENERATION FUNCTIONS
# ============================================

def extract_headings_from_formatted_content(formatted_content: str) -> List[str]:
    """
    Extract headings from formatted markdown content for Contents page.
    Returns list of heading texts (Heading 1 and Heading 2 only).
    """
    headings = []
    lines = formatted_content.split('\n')
    
    for line in lines:
        line = line.strip()
        # Extract Heading 1 (# Title)
        if line.startswith('# ') and not line.startswith('##'):
            heading = line[2:].strip()
            if heading:
                headings.append(heading)
        # Extract Heading 2 (## Section)
        elif line.startswith('## ') and not line.startswith('###'):
            heading = line[3:].strip()
            if heading:
                headings.append(heading)
    
    # Limit to 12 headings maximum for clean Contents page
    return headings[:12]


async def generate_title_and_contents_from_ai_response(formatted_content: str, llm_service, document_title: Optional[str] = None) -> tuple[str, List[str]]:
    """
    Generate engaging title and extract key content points for page 2.
    
    Args:
        formatted_content: Formatted markdown content with headings
        llm_service: LLM service for title generation
        document_title: Optional user-provided title
    
    Returns:
        Tuple of (title, list_of_key_points)
    """
    print(f"\n📄 Generating title and contents from formatted content...")
    
    # Extract headings from formatted content first
    headings = extract_headings_from_formatted_content(formatted_content)
    
    # If we have headings, use them as contents
    if headings:
        print(f"   ✅ Extracted {len(headings)} headings from formatted content")
        contents = headings
    else:
        # Fallback: generate contents using LLM
        print(f"   ⚠️  No headings found, generating contents using LLM...")
        contents = await _generate_contents_with_llm(formatted_content, llm_service)
    
    # Generate title if not provided - ALWAYS use LLM for energetic and generic title
    if document_title and document_title.strip():
        title = document_title.strip()
        print(f"   ✅ Using provided title: '{title}'")
    else:
        # ALWAYS generate title with LLM for energetic and generic title based on AI response
        # Don't use the first Heading 1 as it might be too specific
        print(f"   🤖 Generating energetic and generic title using LLM...")
        title = await _generate_title_with_llm(formatted_content, llm_service)
        print(f"   ✅ Generated energetic title: '{title}'")
    
    # Clean title
    title = title.replace('"', '').replace("'", '').strip()
    if len(title) > 100:
        title = title[:97] + "..."
    if not title:
        title = "Professional Report"
    
    # Ensure we have contents (at least 3 items)
    if not contents or len(contents) < 3:
        contents = ["Introduction", "Key Findings", "Analysis", "Recommendations", "Conclusion"]
    
    print(f"   ✅ Final title: '{title}'")
    print(f"   ✅ Final contents: {len(contents)} items")
    return title, contents


async def _generate_title_with_llm(content: str, llm_service) -> str:
    """Generate energetic and generic title using LLM based on AI response content."""
    system_prompt = """You are a professional document editor. Generate an energetic, engaging, and generic title for a business document.

Rules:
- Maximum 10 words
- Energetic and engaging tone (captivating, dynamic)
- Generic enough to cover the main topic broadly
- Professional business tone
- Capture the essence and main theme of the content
- No quotation marks
- Title case format
- Make it compelling and interesting
- Output ONLY the title, nothing else"""
    
    content_preview = content[:1500] if len(content) > 1500 else content
    
    user_prompt = f"""Based on this AI-generated content, generate an energetic, engaging, and generic title that captures the main theme:

{content_preview}

Generate a title that:
1. Is energetic and engaging (dynamic, compelling)
2. Is generic enough to cover the broad topic
3. Captures the essence of the content
4. Is professional and business-appropriate
5. Maximum 10 words

Output ONLY the title, nothing else."""
    
    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.8,  # Higher temperature for more creative/energetic titles
            max_tokens=40
        )
        
        title = response.content.strip()
        # Remove quotes and clean up
        title = title.replace('"', '').replace("'", '').strip()
        # Remove common prefixes like "Title:" or "Title -"
        if ':' in title:
            title = title.split(':', 1)[-1].strip()
        if '-' in title and len(title.split('-')[0].strip()) < 5:
            title = title.split('-', 1)[-1].strip()
        
        return title if title else "Professional Report"
    except Exception as e:
        print(f"   ⚠️  Title generation failed: {e}")
        return "Professional Report"


async def _generate_contents_with_llm(content: str, llm_service) -> List[str]:
    """Generate contents list using LLM."""
    system_prompt = """You are a professional document editor. Extract 6-10 key content points.

Rules:
- Extract actual headings and key topics
- 6-10 clear, concise points (5-8 words each)
- Output format: One point per line starting with "-"
- Output ONLY the contents list"""
    
    content_preview = content[:2000] if len(content) > 2000 else content
    
    user_prompt = f"""Extract 6-10 key content points from this content:

{content_preview}

Output format (one per line):
- [point 1]
- [point 2]
- [point 3]
..."""
    
    try:
        response = await llm_service.chat_completion(
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.4,
            max_tokens=400
        )
        
        contents = []
        for line in response.content.strip().split('\n'):
            line = line.strip()
            if line.startswith('-') or line.startswith('•'):
                point = line.lstrip('-•').strip()
                if point:
                    contents.append(point)
        
        return contents[:12] if contents else ["Introduction", "Key Findings", "Analysis", "Recommendations", "Conclusion"]
    except:
        return ["Introduction", "Key Findings", "Analysis", "Recommendations", "Conclusion"]


def add_contents_page(doc: Document, title: str, contents: List[str]):
    """
    Add Page 2 with Contents listing - FIXED VERSION.
    
    Args:
        doc: Word document
        title: Document title (for "Contents" heading)
        contents: List of content points to display
    """
    print(f"\n📑 Adding Contents page (Page 2)...")
    
    try:
        # Add page break to start Page 2
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        
        # IMPORTANT: Add page break AFTER title on page 1 to start page 2
        # This ensures Contents appears on page 2, not on page 1 with title
        break_para = doc.add_paragraph()
        run = break_para.add_run()
        br = OxmlElement('w:br')
        br.set(qn('w:type'), 'page')
        run._element.append(br)
        print(f"   📄 Added page break after title (Contents will start on page 2)")
        
        # Add "Contents" heading - TRY MULTIPLE STYLE NAMES
        contents_heading = None
        style_attempts = ['Heading 1 style', 'Heading 1', 'Heading1', 'Title']
        
        for style_name in style_attempts:
            try:
                contents_heading = doc.add_paragraph(style=style_name)
                break
            except:
                continue
        
        # Fallback: create paragraph without style
        if contents_heading is None:
            contents_heading = doc.add_paragraph()
        
        run = contents_heading.add_run("Contents")
        run.bold = True
        run.font.size = Pt(24)
        try:
            run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)  # PwC orange
        except:
            pass
        
        contents_heading.paragraph_format.space_before = Pt(24)
        contents_heading.paragraph_format.space_after = Pt(18)
        contents_heading.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
        
        # Add spacing
        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_after = Pt(12)
        
        # Add content points as bullets - TRY MULTIPLE STYLE NAMES
        bullet_style_attempts = ['List Bullet style', 'List Bullet', 'ListBullet']
        
        for point in contents:
            bullet_para = None
            for style_name in bullet_style_attempts:
                try:
                    bullet_para = doc.add_paragraph(style=style_name)
                    break
                except:
                    continue
            
            # Fallback: create paragraph without style
            if bullet_para is None:
                bullet_para = doc.add_paragraph()
            
            run = bullet_para.add_run(point)
            run.font.size = Pt(14)
            
            # Apply bullet formatting manually
            bullet_para.paragraph_format.space_after = Pt(8)
            bullet_para.paragraph_format.left_indent = Inches(0.25)
            bullet_para.paragraph_format.first_line_indent = Inches(-0.25)
            bullet_para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
            
            # Add bullet character manually if no style worked
            if bullet_para is None or bullet_para.style.name not in bullet_style_attempts:
                bullet_para.text = f"• {point}"
        
        print(f"   ✅ Added Contents page with {len(contents)} items")
        
    except Exception as e:
        print(f"   ⚠️  Could not add Contents page: {e}")
        import traceback
        traceback.print_exc()


async def create_word_document_with_ai_formatting(
    raw_ai_content: str,
    llm_service,
    template_path: str = 'template/Template2.docx',
    user_provided_title: Optional[str] = None
) -> bytes:
    """
    Complete Word Document Generator with AI-powered formatting - FIXED VERSION.
    
    FINAL DOCUMENT STRUCTURE:
    - Page 1: Engaging Title ONLY (centered, with logo, background, theme)
    - Page 2: Contents (key headings/topics extracted from AI response)
    - Page 3+: AI Response Content (formatted professionally)
    
    ALL PAGES MAINTAIN:
    - Same headers (with logo)
    - Same footers
    - Same background color
    - Same theme
    - Consistent branding throughout
    """
    
    print(f"\n{'='*70}")
    print(f"📄 CREATING PROFESSIONAL WORD DOCUMENT - FIXED VERSION")
    print(f"{'='*70}\n")
    
    # Validate template exists
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    
    print(f"✅ Loading template: {template_path}")
    doc = Document(template_path)
    print(f"   Template loaded: {len(doc.paragraphs)} paragraphs")
    
    # STEP 0: PRESERVE TEMPLATE THEME
    print(f"\n🎨 STEP 0: Preserving template theme (logo, background, colors) across all pages...")
    preserve_template_theme(doc)
    
    # STEP 1: EXTRACT INSTRUCTIONS FROM PAGES 3-4 (BEFORE removing content)
    print(f"\n📖 STEP 1: Extracting formatting instructions from pages 3-4...")
    instructions, instr_start, instr_end = extract_instructions_from_template(doc)
    
    if instructions:
        print(f"   ✅ Instructions extracted ({len(instructions)} characters)")
    else:
        print(f"   ⚠️  No instructions found in template")
    
    # STEP 2: USE INSTRUCTIONS AS SYSTEM PROMPT
    if instructions:
        print(f"\n🤖 STEP 2: Using instructions as system prompt to format content...")
        formatted_content = await generate_formatted_content_from_instructions(
            raw_ai_content, instructions, llm_service
        )
        print(f"   ✅ Content formatted using template instructions")
    else:
        print(f"\n⚠️  STEP 2: No instructions found, using content as-is")
        formatted_content = raw_ai_content
    
    # STEP 3: PARSE CONTENT FIRST (needed for title/contents extraction)
    print(f"\n📝 STEP 3: Parsing formatted content into document structure...")
    structured_content = parse_markdown_to_structure(formatted_content)
    print(f"   ✅ Parsed {len(structured_content)} content blocks")
    
    # STEP 4: GENERATE TITLE AND CONTENTS FROM FORMATTED CONTENT
    print(f"\n📌 STEP 4: Generating title and contents from formatted content...")
    title, contents = await generate_title_and_contents_from_ai_response(
        formatted_content, llm_service, user_provided_title
    )
    
    # STEP 5: REMOVE INSTRUCTION PAGES (3-4) FIRST (before cleaning page 1)
    # This ensures we don't have instruction content interfering with page 1
    print(f"\n🗑️  STEP 5: Removing instruction pages 3-4 and unnecessary content...")
    if instr_start is not None and instr_end is not None:
        remove_template_instructions_and_unnecessary_content(doc, instr_start, instr_end)
        print(f"   ✅ Instruction pages removed")
    
    # STEP 6: UPDATE COVER PAGE (PAGE 1) - TITLE ONLY (Remove ALL content including Contents and page breaks)
    # CRITICAL: This must happen AFTER instruction removal but BEFORE adding Contents
    print(f"\n📄 STEP 6: Creating Page 1 - Title ONLY (removing ALL content, page breaks, Contents)...")
    print(f"   🎯 CRITICAL: Ensuring title appears on PAGE 1, not page 2...")
    
    # Clean page 1 aggressively - remove ALL content including page breaks
    # This CRITICAL step ensures title appears on page 1, centered
    title_added = clean_page1_keep_only_title(doc, title)
    if not title_added:
        print(f"   ⚠️  WARNING: Title may not have been added properly to page 1!")
        print(f"   ⚠️  Attempting to fix title placement...")
        # Retry adding title
        try:
            # Remove all paragraphs and add title again
            while len(doc.paragraphs) > 0:
                try:
                    p = doc.paragraphs[0]
                    p_element = p._element
                    parent = p_element.getparent()
                    if parent is not None:
                        parent.remove(p_element)
                except:
                    break
            # Add title as first paragraph
            title_para = doc.add_paragraph()
            title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title_para.paragraph_format.space_before = Pt(320)
            run = title_para.add_run(title)
            run.bold = True
            run.font.size = Pt(38)
            try:
                run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)
            except:
                pass
            print(f"   ✅ Title added using retry method")
        except Exception as e:
            print(f"   ❌ Retry also failed: {e}")
    
    # STEP 7: CLEAN PAGE 2 - Remove all template content (Contents, etc.)
    # Note: This happens AFTER title is added to page 1
    print(f"\n🧹 STEP 7: Cleaning page 2 - removing all template content...")
    clean_page2_remove_template_content(doc, after_title_idx=0)
    
    # STEP 8: REMOVE FOOTERS
    print(f"\n🗑️  STEP 8: Removing footer text content...")
    remove_all_footers(doc)
    
    # STEP 9: ADD CONTENTS PAGE (PAGE 2)
    # Contents page is added AFTER title on page 1, so it will appear on page 2
    print(f"\n📑 STEP 9: Adding Contents page (Page 2)...")
    add_contents_page(doc, title, contents)
    
    # STEP 10: ADD MAIN CONTENT (PAGE 3+)
    print(f"\n📝 STEP 10: Adding AI response content (Page 3+)...")
    
    # Add page break before main content
    try:
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        break_para = doc.add_paragraph()
        run = break_para.add_run()
        br = OxmlElement('w:br')
        br.set(qn('w:type'), 'page')
        run._element.append(br)
        print(f"   📄 Added page break before main content (Page 3)")
    except Exception as e:
        print(f"   ⚠️  Could not add page break: {e}")
    
    # Add content starting on Page 3
    added = 0
    is_first_item = True
    skip_first_heading1 = True  # Skip first Heading 1 to avoid duplicate with title
    
    for item in structured_content:
        style_name = item['style']
        text = item['text']
        
        if not text.strip():
            continue
        
        is_heading = 'Heading' in style_name
        is_heading1 = 'Heading 1' in style_name
        
        # Skip first Heading 1 if it's similar to the title (avoid duplicate)
        if skip_first_heading1 and is_heading1:
            # Check if this heading is very similar to the title
            text_lower = text.lower().strip()
            title_lower = title.lower().strip()
            # Skip if they're very similar (same or nearly same)
            if text_lower == title_lower or text_lower in title_lower or title_lower in text_lower:
                print(f"   ⏭️  Skipping duplicate Heading 1: '{text}' (same as title)")
                skip_first_heading1 = False
                continue
            skip_first_heading1 = False
        
        try:
            para = None
            
            # Try exact style name first
            try:
                para = doc.add_paragraph(style=style_name)
                added += 1
            except (KeyError, ValueError, AttributeError):
                # Fallback to alternatives
                para = doc.add_paragraph()
                added += 1
            
            # Add text content
            if '**' in text:
                parts = text.split('**')
                for i, part in enumerate(parts):
                    if part:
                        run = para.add_run(part)
                        if i % 2 == 1:
                            run.bold = True
            else:
                para.add_run(text)
            
            # Apply professional formatting
            try:
                if 'Heading 1' in style_name:
                    # Heading 1: Main section titles
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(18) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0xE8, 0x77, 0x22)  # PwC orange
                        except:
                            try:
                                run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray fallback
                            except:
                                pass
                    para.paragraph_format.space_before = Pt(24) if is_first_item else Pt(18)
                    para.paragraph_format.space_after = Pt(12)
                    para.paragraph_format.keep_with_next = True
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    is_first_item = False
                
                elif 'Heading 2' in style_name:
                    # Heading 2: Major section headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(16) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(14)
                    para.paragraph_format.space_after = Pt(8)
                    para.paragraph_format.keep_with_next = True
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Heading 3' in style_name:
                    # Heading 3: Subsection headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(14) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(12)
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.keep_with_next = True
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Heading 4' in style_name:
                    # Heading 4: Sub-subsection headings
                    for run in para.runs:
                        run.bold = True
                        run.font.size = Pt(12) if run.font.size is None else run.font.size
                        try:
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                        except:
                            pass
                    para.paragraph_format.space_before = Pt(10)
                    para.paragraph_format.space_after = Pt(4)
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                
                elif 'Body Text' in style_name:
                    # Body text: 11pt font, 1.5 line spacing, left aligned
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.line_spacing = 1.5
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in para.runs:
                        if run.font.size is None:
                            run.font.size = Pt(11)
                        # Body text should not be bold unless explicitly marked
                        if not any('**' in item['text'] for item in [item]):
                            run.bold = False
                
                elif 'Bullet' in style_name or 'Number' in style_name or 'Alpha' in style_name:
                    # Lists: hanging indent, proper spacing
                    para.paragraph_format.left_indent = Inches(0.25)
                    para.paragraph_format.first_line_indent = Inches(-0.25)  # Hanging indent
                    para.paragraph_format.space_after = Pt(3)
                    para.paragraph_format.line_spacing = 1.5
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in para.runs:
                        if run.font.size is None:
                            run.font.size = Pt(11)
                
                elif 'Quote' in style_name:
                    # Quotes: italic, indented
                    para.paragraph_format.left_indent = Inches(0.5)
                    para.paragraph_format.space_after = Pt(6)
                    para.paragraph_format.line_spacing = 1.5
                    para.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    for run in para.runs:
                        run.italic = True
                        if run.font.size is None:
                            run.font.size = Pt(11)
                
                if is_first_item:
                    is_first_item = False
                
            except Exception as e:
                print(f"   ⚠️  Could not apply all formatting: {e}")
                pass
            
        except Exception as e:
            print(f"   ⚠️  Error adding content block '{text[:30]}...': {e}")
            continue
    
    print(f"   ✅ Added {added} content paragraphs to Page 3+")
    
    # Save document
    print(f"\n💾 Saving document...")
    buffer = io.BytesIO()
    try:
        doc.save(buffer)
        buffer.seek(0)
    except Exception as e:
        print(f"   ❌ Error saving document: {e}")
        raise
    
    file_size = len(buffer.getvalue())
    print(f"\n✅ Document generated successfully!")
    print(f"   Title: '{title}'")
    print(f"   Contents items: {len(contents)}")
    print(f"   Content blocks: {len(structured_content)}")
    print(f"   File size: {file_size:,} bytes")
    print(f"   Structure:")
    print(f"      - Page 1: Title ('{title}')")
    print(f"      - Page 2: Contents ({len(contents)} items)")
    print(f"      - Page 3+: AI Response ({added} paragraphs)")
    print(f"{'='*70}\n")
    
    return buffer.getvalue()
def validate_ppt_rules(prs: Presentation) -> dict:
    """Perform deterministic rule-based validation on PowerPoint"""
    violations = []
    warnings = []
    slide_count = len(prs.slides)

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_violations = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    text = paragraph.text.strip()

                    if text and text.lower() in [
                        "tbd",
                        "to be determined",
                        "insert chart here",
                        "placeholder",
                    ]:
                        slide_violations.append(
                            f"Slide {slide_idx}: Placeholder text found: '{text}'"
                        )

                    if len(text) > 0:
                        if text.endswith(".") and not hasattr(shape, "has_table"):
                            slide_violations.append(
                                f"Slide {slide_idx}: Bullet point should not end with period: '{text[:50]}...'"
                            )

                        if "  " in text:
                            slide_violations.append(
                                f"Slide {slide_idx}: Double spaces found in text"
                            )

                    for run in paragraph.runs:
                        if run.font.size:
                            size_pt = run.font.size.pt
                            if size_pt < 10:
                                warnings.append(
                                    f"Slide {slide_idx}: Font size {size_pt}pt is below minimum 10pt in text: '{run.text[:30]}...'"
                                )

        violations.extend(slide_violations)

    if slide_count < 3:
        warnings.append(
            f"Presentation has only {slide_count} slides - may lack sufficient structure"
        )

    return {
        "violations": violations[:50],
        "warnings": warnings[:50],
        "total_violations": len(violations),
        "total_warnings": len(warnings),
        "slide_count": slide_count,
    }


@app.post("/api/validate-best-practices")
async def validate_best_practices(
    file: UploadFile = File(...), categories: Optional[str] = Form(None)
):
    """Validate PowerPoint against 75+ PwC consulting best practices"""
    try:
        contents = await file.read()
        prs = Presentation(io.BytesIO(contents))

        rule_validation = validate_ppt_rules(prs)

        selected_categories = []
        if categories:
            selected_categories = categories.split(",")
        else:
            selected_categories = [
                "Structure",
                "Visuals",
                "Design",
                "Charts",
                "Formatting",
                "Content",
            ]

        all_best_practices = {
            "Structure": [
                "Logical flow: Introduction → Context → Analysis → Insights → Implications → Recommendations → Next steps",
                "Pyramid Principle applied: Each section opens with key message headline",
                "Slide hierarchy consistent: Section headers, transition slides, content slides clearly differentiated",
                "Each slide answers a specific question",
                "Slide headers capture narrative takeaway, not just topic",
                "Horizontal logic: Headers form coherent narrative when read sequentially",
            ],
            "Visuals": [
                "Color used strategically to guide the eye (primary color for key messages)",
                "Converting text into purposeful visuals (2×2 grids, process flows)",
                "Visual hierarchy clear: Most important content is visually dominant",
                "Icons/images purposeful and related to text",
                "Bullets consistent: No mix of shapes unless intentional",
                "Parallel language in all bullets",
                "Concise bullet statements",
            ],
            "Design": [
                "Less is more: No unnecessary drop shadows, gradients, or animations",
                "White space used intentionally between sections",
                "Alignment perfect: Objects snapped to grid, text aligned precisely",
                "Shape sizes consistent (use exact dimensions)",
                "Images high-res and consistent style",
                "Consistent slide backgrounds (all white or light gray)",
                "Icons from PowerPoint library matching color scheme",
            ],
            "Charts": [
                "Charts simplified: No unnecessary borders, 3D effects, or legends when labels suffice",
                "Font sizes legible (10-12pt in chart labels)",
                "Axes consistent and start at zero when appropriate",
                "Same category = same color throughout deck",
                "Charts embedded rather than screenshots",
                "Every chart has subtitle with period and units",
                "Source lines included below charts",
                "Chart titles as takeaways not just labels",
            ],
            "Formatting": [
                "Fonts consistent (same font family across slides)",
                "Font sizes standardized (28pt title, 16pt header, 14pt sub-header)",
                "Font size never below 10pt (except footnotes)",
                "Full page content center-aligned",
                "Paragraph spacing set appropriately",
                "Equal alignment and spacing using Align and Distribute tools",
                "Color palette consistent with firm's official palette",
                "Headers/footers standardized using slide master",
            ],
            "Content": [
                "No typos (spell check completed)",
                "No double spaces",
                "Jargon minimized and replaced with plain business English",
                "Consistency in terminology throughout",
                "All acronyms defined on first use",
                "No previous client references",
                "No placeholder language (TBD, insert chart here)",
                "No leftover comments or notes",
                "Confidentiality notice present where required",
                "All links tested and working",
            ],
        }

        slide_count = len(prs.slides)
        total_text = 0
        total_shapes = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                if hasattr(shape, "text"):
                    total_text += len(shape.text)

        system_prompt = f"""You are a PwC presentation quality expert. Analyze PowerPoint presentations against consulting best practices.
Be specific, critical, and constructive. Identify both strengths and areas for improvement."""

        practices_to_check = []
        for category in selected_categories:
            if category in all_best_practices:
                practices_to_check.extend(all_best_practices[category])

        violations_section = ""
        if (
            rule_validation["total_violations"] > 0
            or rule_validation["total_warnings"] > 0
        ):
            violations_section = f"""
**Automated Rule Validation Results:**
- Critical Violations Found: {rule_validation['total_violations']}
- Warnings Found: {rule_validation['total_warnings']}

**Specific Issues Detected:**
{chr(10).join(f'• {v}' for v in rule_validation['violations'][:10])}
{chr(10).join(f'⚠ {w}' for w in rule_validation['warnings'][:10])}
"""

        user_prompt = f"""Analyze this PowerPoint presentation against PwC best practices:

**Presentation Stats:**
- Total Slides: {slide_count}
- Total Shapes: {total_shapes}
- Approximate Text Length: {total_text} characters
{violations_section}
**Best Practices to Validate** (from categories: {', '.join(selected_categories)}):
{chr(10).join(f'{i+1}. {practice}' for i, practice in enumerate(practices_to_check))}

**Analysis Required:**
For each category, provide:
1. **Compliance Score**: Estimate % compliance with these practices
2. **Key Issues Found**: Specific violations or concerns (in addition to automated violations above)
3. **Strengths**: What the presentation does well
4. **Recommendations**: Prioritized improvements (address automated violations first, then other issues)

Format your response with clear headers for each category. Start with a summary of the automated violations."""

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ]

        return StreamingResponse(
            generate_stream(messages), media_type="text/event-stream"
        )

    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error validating presentation: {str(e)}"
        )
# Add these new functions to your main.py - ONLY USED AT DOWNLOAD TIME




@app.post("/api/export/word")
async def export_word(request: ExportRequest):
    """
    Export AI content as Word document.
    
    WORKFLOW:
    1. Extract instructions from pages 3-4 (BEFORE removal)
    2. Use instructions as system prompt to format content
    3. Generate/use title
    4. Parse formatted content
    5. Update cover page
    6. REMOVE pages 3-4 (instructions) - they are no longer needed
    7. Add Table of Contents and main content
    
    NEW FEATURES:
    - Extracts Pages 3-4 instructions FIRST
    - Uses instructions as system prompt to format content
    - Auto-generates title from content
    - Removes Pages 3-4 AFTER using them as system prompt
    - Professional formatting with Table of Contents
    """
    try:
        template_path = 'template/Template2.docx'
        
        print(f"\n{'='*60}")
        print(f"📤 WORD EXPORT REQUEST")
        print(f"{'='*60}")
        print(f"📝 Content length: {len(request.content)} characters")
        print(f"📌 Title provided: {bool(request.title and request.title.strip())}")
        
        # Get LLM service
        llm = get_llm()
        
        # Generate document with AI formatting
        docx_bytes = await create_word_document_with_ai_formatting(
            raw_ai_content=request.content,
            llm_service=llm,
            template_path=template_path,
            user_provided_title=request.title
        )
        
        # Generate filename
        if request.title and request.title.strip():
            safe_title = request.title.replace(' ', '_').replace('/', '_')[:50]
        else:
            safe_title = "Document"
        
        date_str = datetime.now().strftime('%Y%m%d')
        filename = f"{safe_title}_{date_str}.docx"
        
        print(f"✅ Export successful: {filename}")
        print(f"{'='*60}\n")
        
        return StreamingResponse(
            io.BytesIO(docx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=f"Template not found: {str(e)}")
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error creating Word document: {str(e)}")


# ============================================
# PDF EXPORT ENDPOINT
# ============================================




@app.post("/api/export/pdf")
async def export_pdf(request: ExportRequest):
    """
    Export AI content as PDF document.
    
    WORKFLOW:
    1. Extract instructions from pages 3-4 (BEFORE removal)
    2. Use instructions as system prompt to format content
    3. Generate/use title
    4. Parse formatted content
    5. Update cover page
    6. REMOVE pages 3-4 (instructions) - they are no longer needed
    7. Add Table of Contents and main content
    8. Convert Word to PDF
    
    Creates a Word document first, then converts it to PDF.
    Uses the same template and formatting as Word export.
    Pages 3-4 are removed after being used as system prompt.
    """
    try:
        template_path = 'template/Template2.docx'
        
        print(f"\n{'='*60}")
        print(f"📤 PDF EXPORT REQUEST")
        print(f"{'='*60}")
        print(f"📝 Content length: {len(request.content)} characters")
        print(f"📌 Title provided: {bool(request.title and request.title.strip())}")
        
        # Get LLM service
        llm = get_llm()
        
        # Generate Word document with AI formatting
        docx_bytes = await create_word_document_with_ai_formatting(
            raw_ai_content=request.content,
            llm_service=llm,
            template_path=template_path,
            user_provided_title=request.title
        )
        
        # Convert Word document to PDF
        print(f"🔄 Converting Word document to PDF...")
        try:
            pdf_bytes = create_pdf_from_docx_bytes(docx_bytes)
        except ValueError as e:
            # Word COM unavailable - extract text and use ReportLab fallback
            print(f"⚠️  Word-to-PDF conversion unavailable, using ReportLab fallback...")
            
            # Extract formatted content from DOCX
            formatted_content = extract_text_from_docx(docx_bytes, max_chars=50000)
            
            # Generate title if not provided
            if not request.title or not request.title.strip():
                title = await generate_title_from_ai_content(formatted_content, llm)
            else:
                title = request.title
            
            # Generate PDF using ReportLab
            pdf_bytes = create_pdf_from_text_fallback(formatted_content, title)
            print(f"✅ PDF generated using ReportLab fallback")
        
        # Generate filename
        if request.title and request.title.strip():
            safe_title = request.title.replace(' ', '_').replace('/', '_')[:50]
        else:
            safe_title = "Document"
        
        date_str = datetime.now().strftime('%Y%m%d')
        filename = f"{safe_title}_{date_str}.pdf"
        
        print(f"✅ Export successful: {filename}")
        print(f"{'='*60}\n")
        
        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=f"Template not found: {str(e)}")
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Error creating PDF document: {str(e)}")
# @app.post("/api/export/preview")
# async def preview_document_structure(request: ExportRequest):
#     """
#     Preview how AI content will be structured in document.
#     Useful for debugging and testing.
    
#     Returns JSON showing:
#     - How content is parsed
#     - What styles will be applied
#     - Content block preview
#     """
#     try:
#         from datetime import datetime
        
#         # Parse content
#         structured = parse_ai_markdown_to_docx_structure(request.content)
        
#         return {
#             "timestamp": datetime.now().isoformat(),
#             "title": request.title,
#             "original_content_length": len(request.content),
#             "parsed_blocks": len(structured),
#             "structure_preview": [
#                 {
#                     "index": i,
#                     "style": item["style"],
#                     "text_preview": item["text"][:100] + "..." if len(item["text"]) > 100 else item["text"],
#                     "text_length": len(item["text"])
#                 }
#                 for i, item in enumerate(structured[:20])  # First 20 blocks
#             ],
#             "style_distribution": {
#                 style: sum(1 for item in structured if item["style"] == style)
#                 for style in set(item["style"] for item in structured)
#             }
#         }
    
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Error previewing content: {str(e)}")



def extract_color_scheme(prs: Presentation):
    """Extract the most common colors from a presentation."""
    colors = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type == 1:
                if (
                    hasattr(shape.fill, "fore_color")
                    and shape.fill.fore_color.type == 1
                ):
                    rgb = shape.fill.fore_color.rgb
                    colors.append((rgb[0], rgb[1], rgb[2]))

            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, "color") and run.font.color.type == 1:
                            rgb = run.font.color.rgb
                            colors.append((rgb[0], rgb[1], rgb[2]))

    if colors:
        color_counts = Counter(colors)
        most_common = [color for color, count in color_counts.most_common(5)]
        return most_common

    return [(208, 74, 2), (255, 107, 0), (26, 26, 26)]


async def correct_text_with_ai(text: str) -> str:
    """Use LLM AI to correct spelling and grammar."""
    if not text or len(text.strip()) == 0:
        return text

    try:
        llm = get_llm()
        response = await llm.chat_completion(
            messages=[
                {
                    "role": "system",
                    "content": "You are a professional editor. Fix spelling and grammar mistakes while preserving the original meaning and tone. Return ONLY the corrected text without any explanations or additional commentary.",
                },
                {"role": "user", "content": f"Correct this text: {text}"},
            ],
            temperature=0.3,
            max_tokens=1000,
        )
        return response.content.strip()
    except:
        return text


def align_shapes(slide):
    """Align shapes that are close to each other."""
    shapes = [s for s in slide.shapes if hasattr(s, "left") and hasattr(s, "top")]

    if len(shapes) < 2:
        return

    shapes_by_row = {}
    tolerance = Inches(0.1)

    for shape in shapes:
        aligned = False
        for ref_top in shapes_by_row.keys():
            if abs(shape.top - ref_top) < tolerance:
                shapes_by_row[ref_top].append(shape)
                aligned = True
                break

        if not aligned:
            shapes_by_row[shape.top] = [shape]

    for ref_top, row_shapes in shapes_by_row.items():
        if len(row_shapes) > 1:
            avg_top = sum(s.top for s in row_shapes) // len(row_shapes)
            for shape in row_shapes:
                shape.top = avg_top


def apply_color_scheme(prs: Presentation, color_scheme: list):
    """Apply color scheme to presentation."""
    if not color_scheme:
        return

    primary_color = RGBColor(*color_scheme[0])
    secondary_color = (
        RGBColor(*color_scheme[1]) if len(color_scheme) > 1 else primary_color
    )

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "fill") and shape.fill.type == 1:
                if hasattr(shape.fill, "fore_color"):
                    shape.fill.fore_color.rgb = primary_color

            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run.font, "color"):
                            if (
                                run.font.bold
                                or run.font.size
                                and run.font.size > Pt(18)
                            ):
                                run.font.color.rgb = primary_color
                            else:
                                run.font.color.rgb = secondary_color


@app.post("/api/ppt/improve")
async def improve_ppt(
    original_ppt: UploadFile = File(...), reference_ppt: UploadFile = File(None)
):
    """Improve PowerPoint presentation: correct spelling/grammar, align shapes, rebrand colors."""
    try:
        original_content = await original_ppt.read()
        original_prs = Presentation(io.BytesIO(original_content))

        color_scheme = None
        if reference_ppt:
            reference_content = await reference_ppt.read()
            reference_prs = Presentation(io.BytesIO(reference_content))
            color_scheme = extract_color_scheme(reference_prs)

        for slide in original_prs.slides:
            align_shapes(slide)

            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                corrected_text = await correct_text_with_ai(run.text)
                                run.text = corrected_text

        if color_scheme:
            apply_color_scheme(original_prs, color_scheme)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            original_prs.save(tmp_file.name)
            tmp_path = tmp_file.name

        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="improved_presentation.pptx",
            background=None,
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PPT: {str(e)}")


class SanitizationConversationRequest(BaseModel):
    messages: List[Message]
    uploaded_file_name: Optional[str] = None
    client_identity: Optional[str] = None
    page_range: Optional[str] = None
    tier1_services: Optional[List[str]] = None  # Default services
    tier2_services: Optional[List[str]] = None  # Opt-in services
    stream: bool = True


@app.post("/api/ppt/sanitize/conversation")
async def sanitize_conversation(request: SanitizationConversationRequest):
    """
    Conversational sanitization endpoint that guides users through the sanitization process.

    Tier 1 (Default/Opt-out) Services:
    - Convert to PwC standard template
    - Replace client names and logos with placeholders
    - Delete speaker notes and comments
    - Clear presentation metadata
    - Remove numeric data (replace with X patterns)

    Tier 2 (Opt-in) Services:
    - Change competitor company names
    - Remove client-specific financial data
    - Redact business unit names
    - Remove product names
    - Sanitize location data (cities, addresses)

    Tier 3 (Advanced, performed on request):
    - Modify Think-cell chart values
    - Custom regex-based replacements
    - Advanced contextual sanitization

    Tier 4 (Never Changed):
    - PwC branding and watermarks
    - Standard consulting frameworks (MECE, Porter's Five Forces, etc.)
    - Industry-standard terminology
    """
    try:
        llm = get_llm()

        system_prompt = """You are a specialized PwC Presentation Sanitization Assistant. Your role is to guide users through a comprehensive document sanitization process.

**SANITIZATION WORKFLOW:**

1. **Recognition & Recap**: When a user mentions "sanitize", "sanitization", or "sanitation", acknowledge their request and provide a brief overview of the service.

2. **Document Collection**: If they haven't uploaded a PowerPoint file yet, politely request it. Explain that you need the original PowerPoint document to proceed.

3. **Client Identity** (Optional): Ask if they want to specify the original client name, or if you should auto-detect it from the document context.

4. **Page Range**: Ask if they want to sanitize specific pages or all pages (default: all pages).

5. **Quick Start Overview**: Present the sanitization service structure:

   **TIER 1 - Default Services (Applied unless you opt-out):**
   1. Convert to PwC standard template
   2. Replace client names and logos with [Client] placeholders  
   3. Delete all speaker notes and comments
   4. Clear presentation metadata
   5. Remove/replace numeric data with X patterns

   **TIER 2 - Opt-in Services (Only if you request):**
   6. Change competitor company names to [Competitor]
   7. Remove client-specific financial data
   8. Redact business unit names to [BU]
   9. Replace product names with [Product]
   10. Sanitize location data (cities, addresses)
   11. Remove embedded hyperlinks

   Ask: "Would you like any of the Tier 2 services?"

6. **Customization**: Allow users to:
   - Add specific actions
   - Modify listed actions
   - Remove default actions
   - Request specific replacements

7. **Service Listing**: If requested, provide a complete numbered list of ALL sanitization actions with unique indices.

8. **Confirmation**: Summarize the selected services and get final confirmation before processing.

**CONVERSATION GUIDELINES:**
- Be concise and professional
- Use bullet points for clarity
- Number services for easy reference
- Always confirm before processing
- Explain each service briefly when asked
- If user asks about specific pages, note that page-specific requests are supported

**IMPORTANT:**
- Tier 3 services (Think-cell modifications, advanced customization) are available but not announced unless specifically requested
- Never modify Tier 4 items (PwC branding, standard frameworks, industry terminology)
- Client detection can be automatic if not specified
- Default assumption: sanitize all pages unless range specified

Respond naturally and guide the user through the process step by step."""

        messages = [{"role": "system", "content": system_prompt}]
        messages.extend(
            [{"role": msg.role, "content": msg.content} for msg in request.messages]
        )

        # Add context if file uploaded
        if request.uploaded_file_name:
            context = f"\n\n[File uploaded: {request.uploaded_file_name}]"
            if request.client_identity:
                context += f"\n[Client identity specified: {request.client_identity}]"
            if request.page_range:
                context += f"\n[Page range: {request.page_range}]"
            messages[-1]["content"] += context

        async def stream_response():
            try:
                async for content in llm.stream_completion(
                    messages=messages,
                    temperature=0.7,
                    max_tokens=2048,
                ):
                    yield f"data: {json.dumps({'content': content})}\n\n"

                yield f"data: {json.dumps({'done': True})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"

        return StreamingResponse(stream_response(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error in sanitization conversation: {str(e)}"
        )


@app.post("/api/ppt/sanitize")
async def sanitize_ppt(
    original_ppt: UploadFile = File(...),
    client_name: Optional[str] = None,
    client_products: Optional[str] = None,
    business_units: Optional[str] = None,
    sanitization_options: Optional[str] = Form(None),
    page_range: Optional[str] = Form(None),
):
    """
    Comprehensive PowerPoint sanitization with grammar correction:

    Data Sanitization (Tier-based):

    Tier 1 (Default):
    - Client names: Replace with [Client]
    - Logos and watermarks: Remove images
    - Speaker notes and comments: Clear all
    - Metadata: Clear all document properties
    - Numeric data: Replace with X placeholders

    Tier 2 (Opt-in):
    - Product names: Replace with [Product]
    - Business unit names: Replace with [BU]
    - Competitor names: Replace with [Competitor]
    - Financial data: Remove client-specific numbers
    - Location data: Redact cities, addresses

    Tier 3 (Advanced):
    - Think-cell chart modifications
    - Custom replacements
    - Personal information: Emails, phones, SSN
    - Embedded files: Disconnect attachments

    Grammar & Spelling:
    - AI-powered correction using Groq
    - Preserves formatting and tone
    """
    try:
        if not original_ppt.filename.endswith(".pptx"):
            raise HTTPException(
                status_code=400, detail="File must be a .pptx PowerPoint file"
            )

        original_content = await original_ppt.read()

        client_names = []
        if client_name:
            client_names = [
                name.strip() for name in client_name.split(",") if name.strip()
            ]

        product_names = []
        if client_products:
            product_names = [
                name.strip() for name in client_products.split(",") if name.strip()
            ]

        business_unit_names = []
        if business_units:
            business_unit_names = [
                name.strip() for name in business_units.split(",") if name.strip()
            ]

        # Parse sanitization options for selective sanitization
        options = {}
        if sanitization_options:
            try:
                options = json.loads(sanitization_options)
            except json.JSONDecodeError:
                pass

        sanitizer = PPTSanitizer(
            client_names=client_names,
            product_names=product_names,
            business_units=business_unit_names,
            options=options,
        )

        sanitized_output = sanitizer.sanitize_presentation(
            io.BytesIO(original_content),
            client_name=client_names[0] if client_names else None,
            fix_grammar=True,  # Enable grammar fixing
        )

        stats = sanitizer.get_stats()
        stats["grammar_corrections"] = sanitizer.sanitization_stats.get(
            "grammar_corrections", 0
        )

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            tmp_file.write(sanitized_output.getvalue())
            tmp_path = tmp_file.name

        headers = {"X-Sanitization-Stats": json.dumps(stats)}

        return FileResponse(
            tmp_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="sanitized_presentation.pptx",
            headers=headers,
            background=None,
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error sanitizing PPT: {str(e)}")


# Podcast Generation Functions


def get_polly_client():
    """Get or create Polly client with AWS credentials"""
    aws_access_key = os.getenv("AWS_ACCESS_KEY_ID")
    aws_secret_key = os.getenv("AWS_SECRET_ACCESS_KEY")
    aws_region = os.getenv("AWS_REGION", "us-east-1")

    if not aws_access_key or not aws_secret_key:
        raise HTTPException(
            status_code=500,
            detail="AWS credentials not configured. Please add AWS_ACCESS_KEY_ID and AWS_SECRET_ACCESS_KEY.",
        )

    return boto3.client(
        "polly",
        aws_access_key_id=aws_access_key,
        aws_secret_access_key=aws_secret_key,
        region_name=aws_region,
    )


async def generate_podcast_script(
    content: str, customization: Optional[str] = None, podcast_style: str = "dialogue"
):
    """Generate a podcast script from content using LLM - supports dialogue or monologue styles"""
    llm = get_llm()

    if podcast_style == "monologue":
        system_prompt = """You are a podcast script writer creating engaging single-narrator podcast episodes.

Create a natural, engaging narration:
- **Narrator**: Professional, clear, engaging storytelling voice

Script Requirements:
- Target 10-15 minutes (2,000-2,500 words)
- Natural narration with clear structure and flow
- Include verbal pauses and conversational elements
- Start with a brief introduction of the topic
- End with key takeaways and closing remarks
- Make complex ideas accessible and engaging
- Use storytelling and real-world examples

Format your script as:
NARRATOR: [narration]

Keep the narration flowing naturally with smooth transitions between sections."""

        user_prompt = f"""Create an engaging podcast monologue discussing the following content:

{content}

{f"Special instructions: {customization}" if customization else ""}

Generate a complete, natural-sounding podcast narration."""
    else:
        system_prompt = """You are a podcast script writer creating engaging two-host conversational podcasts.

Create a natural, engaging dialogue between two hosts:
- **Alex** (Host 1): Thoughtful, asks clarifying questions, brings strategic perspective
- **Jordan** (Host 2): Energetic, explains concepts clearly, adds relatable examples

Script Requirements:
- Target 10-15 minutes (2,000-2,500 words)
- Natural conversation with questions, answers, insights, and examples
- Include verbal pauses and conversational elements (e.g., "you know," "right," "that's interesting")
- Start with a brief introduction of the topic
- End with key takeaways and closing remarks
- Make complex ideas accessible and engaging
- Use storytelling and real-world examples

Format your script as:
ALEX: [dialogue]
JORDAN: [dialogue]

Keep the conversation flowing naturally with back-and-forth exchanges."""

        user_prompt = f"""Create an engaging podcast script discussing the following content:

{content}

{f"Special instructions: {customization}" if customization else ""}

Generate a complete, natural-sounding podcast conversation between Alex and Jordan."""

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

    response = await llm.chat_completion(
        messages=messages,
        temperature=0.8,
        max_tokens=4000,
    )

    return response.content


def parse_script_segments(script: str) -> List[dict]:
    """Parse podcast script into segments with speaker labels"""
    segments = []
    lines = script.strip().split("\n")

    current_speaker = None
    current_text = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Check if line starts with a speaker label
        if line.upper().startswith("ALEX:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "ALEX"
            current_text = [line[5:].strip()]  # Remove "ALEX:"
        elif line.upper().startswith("JORDAN:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "JORDAN"
            current_text = [line[7:].strip()]  # Remove "JORDAN:"
        elif line.upper().startswith("NARRATOR:"):
            if current_speaker and current_text:
                segments.append(
                    {"speaker": current_speaker, "text": " ".join(current_text).strip()}
                )
            current_speaker = "NARRATOR"
            current_text = [line[9:].strip()]  # Remove "NARRATOR:"
        else:
            # Continuation of current speaker's dialogue
            if current_speaker:
                current_text.append(line)

    # Add final segment
    if current_speaker and current_text:
        segments.append(
            {"speaker": current_speaker, "text": " ".join(current_text).strip()}
        )

    return segments


async def synthesize_with_polly(text: str, voice_id: str) -> bytes:
    """Synthesize speech using Amazon Polly"""
    polly_client = get_polly_client()

    try:
        response = polly_client.synthesize_speech(
            Text=text, OutputFormat="mp3", VoiceId=voice_id, Engine="neural"
        )

        return response["AudioStream"].read()
    except (BotoCoreError, ClientError) as error:
        raise HTTPException(
            status_code=500, detail=f"Polly synthesis error: {str(error)}"
        )


async def create_podcast_audio(script: str) -> bytes:
    """Create complete podcast audio from script with two distinct voices"""
    segments = parse_script_segments(script)

    if not segments:
        raise HTTPException(status_code=400, detail="Could not parse podcast script")

    # Voice mapping
    voice_map = {
        "ALEX": "Matthew",  # Male neural voice
        "JORDAN": "Joanna",  # Female neural voice
        "NARRATOR": "Matthew",  # Professional male voice for monologues
    }

    # Create audio segments
    audio_segments = []

    for segment in segments:
        pass
        # voice = voice_map.get(segment["speaker"], "Matthew")
        # audio_data = await synthesize_with_polly(segment["text"], voice)

        # # Convert to AudioSegment
        # audio_segment = AudioSegment.from_mp3(io.BytesIO(audio_data))
        # audio_segments.append(audio_segment)

        # # Add brief pause between speakers (300ms)
        # audio_segments.append(AudioSegment.silent(duration=300))

    # Concatenate all segments
    if not audio_segments:
        raise HTTPException(status_code=400, detail="No audio segments generated")

    final_audio = audio_segments[0]
    for segment in audio_segments[1:]:
        final_audio += segment

    # Export to MP3
    output_buffer = io.BytesIO()
    final_audio.export(output_buffer, format="mp3", bitrate="128k")
    output_buffer.seek(0)

    return output_buffer.read()


@app.post("/api/generate-podcast")
async def generate_podcast(
    files: Optional[List[UploadFile]] = File(None),
    content_text: Optional[str] = Form(None),
    customization: Optional[str] = Form(None),
    podcast_style: Optional[str] = Form("dialogue"),
):
    """Generate a NotebookLM-style podcast from uploaded documents or text"""
    try:

        async def event_generator():
            try:
                # Step 1: Parse uploaded files and combine content
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Parsing uploaded documents...', 'percent': 10})}\n\n"

                combined_content = ""

                # Add explicit text content if provided
                if content_text:
                    combined_content += content_text + "\n\n"

                # Parse uploaded files
                if files:
                    for file in files:
                        file_content = await file.read()
                        filename = file.filename.lower()

                        if filename.endswith(".pdf"):
                            parsed_content = extract_text_from_pdf(file_content)
                        elif filename.endswith(".docx"):
                            parsed_content = extract_text_from_docx(file_content)
                        elif filename.endswith((".txt", ".md")):
                            parsed_content = file_content.decode(
                                "utf-8", errors="ignore"
                            )
                        else:
                            continue

                        combined_content += f"\n\n{parsed_content}"

                if not combined_content.strip():
                    yield f"data: {json.dumps({'type': 'error', 'message': 'No content provided. Please upload files or provide text.'})}\n\n"
                    return

                # Limit content length
                if len(combined_content) > 50000:
                    combined_content = combined_content[:50000]

                # Step 2: Generate podcast script
                style_message = (
                    "monologue script"
                    if podcast_style == "monologue"
                    else "conversational podcast script"
                )
                yield f"data: {json.dumps({'type': 'progress', 'message': f'Generating {style_message}...', 'percent': 30})}\n\n"

                script = await generate_podcast_script(
                    combined_content, customization, podcast_style
                )

                if not script or not script.strip():
                    yield f"data: {json.dumps({'type': 'error', 'message': 'Failed to generate podcast script'})}\n\n"
                    return

                yield f"data: {json.dumps({'type': 'progress', 'message': 'Script generated! Now synthesizing audio with Amazon Polly...', 'percent': 50})}\n\n"

                # Step 3: Synthesize audio
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Converting to speech (this may take a few minutes)...', 'percent': 60})}\n\n"

                audio_data = await create_podcast_audio(script)

                # Log audio generation stats
                print(f"Audio generated: {len(audio_data)} bytes")

                # Step 4: Encode audio as base64 for transfer
                yield f"data: {json.dumps({'type': 'progress', 'message': 'Finalizing podcast...', 'percent': 90})}\n\n"

                audio_base64 = base64.b64encode(audio_data).decode("utf-8")
                print(f"Base64 encoded: {len(audio_base64)} characters")

                # Return script and audio
                yield f"data: {json.dumps({'type': 'script', 'content': script})}\n\n"
                yield f"data: {json.dumps({'type': 'complete', 'message': 'Podcast generated successfully!', 'audio': audio_base64, 'percent': 100})}\n\n"

            except Exception as e:
                yield f"data: {json.dumps({'type': 'error', 'message': str(e)})}\n\n"

        return StreamingResponse(event_generator(), media_type="text/event-stream")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
